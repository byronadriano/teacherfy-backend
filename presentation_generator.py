import tempfile
from pptx import Presentation
from flask import send_file

def generate_presentation(outline_text):
    prs = Presentation()
    
    # Parse the outline text to extract slide content
    slides_content = parse_outline_text(outline_text)
    
    for slide_info in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using layout with title and content
        
        # Set the title
        title = slide.shapes.title
        title.text = slide_info['title']
        
        # Set the content
        content = slide.placeholders[1]  # Index 1 is typically the content placeholder
        tf = content.text_frame
        
        # Add main content
        p = tf.add_paragraph()
        p.text = slide_info['content']
        
        # Add bullet points if any
        if 'bullets' in slide_info and slide_info['bullets']:
            for bullet in slide_info['bullets']:
                bullet_p = tf.add_paragraph()
                bullet_p.text = bullet
                bullet_p.level = 1  # Set indentation level for bullets
    
    # Save to temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        tmp.seek(0)
        return tmp.name

def parse_outline_text(outline_text):
    """
    Parse the outline text into structured slide content
    """
    slides = []
    current_slide = None
    
    # Split by lines and process
    lines = outline_text.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check if this is a new slide (starts with "Slide" or contains ":")
        if line.lower().startswith('slide') or ':' in line:
            if current_slide:
                slides.append(current_slide)
            
            # Extract title
            if ':' in line:
                title = line.split(':', 1)[0].strip()
                content = line.split(':', 1)[1].strip()
            else:
                title = line
                content = ""
                
            current_slide = {
                'title': title,
                'content': content,
                'bullets': []
            }
        
        # If line starts with - or *, treat as bullet point
        elif line.strip().startswith(('-', '*', '•')):
            if current_slide:
                bullet_text = line.lstrip('-*• ').strip()
                current_slide['bullets'].append(bullet_text)
        
        # Otherwise, append to current slide's content
        else:
            if current_slide:
                if current_slide['content']:
                    current_slide['content'] += '\n' + line
                else:
                    current_slide['content'] = line
    
    # Add the last slide if exists
    if current_slide:
        slides.append(current_slide)
    
    return slides