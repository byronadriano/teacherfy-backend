# src/resource_handlers/presentation_handler.py - CLEANED VERSION
import os
import logging
from typing import Dict, Any, List
from .base_handler import BaseResourceHandler

logger = logging.getLogger(__name__)

class PresentationHandler(BaseResourceHandler):
    """Handler for generating PowerPoint presentations with optional image support."""
    
    def __init__(self, structured_content: List[Dict[str, Any]], **kwargs):
        super().__init__(structured_content, **kwargs)
        logger.info(f"PresentationHandler initialized with images: {self.include_images}")
    
    def generate(self) -> str:
        """Generate the presentation file and return the file path."""
        # Create temp file
        temp_file = self.create_temp_file("pptx")
        
        # Create presentation with clean structure and image preference
        logger.info(f"Creating presentation with {len(self.structured_content)} slides, images: {self.include_images}")
        
        # Use the built-in presentation creation
        prs = self._create_presentation_from_structured_content(self.structured_content, self.include_images)
        
        # Save presentation
        logger.info(f"Saving presentation to {temp_file}")
        prs.save(temp_file)
        
        # Verify file was created and is not empty
        if not os.path.exists(temp_file):
            raise FileNotFoundError(f"Failed to create presentation file at {temp_file}")
            
        file_size = os.path.getsize(temp_file)
        logger.info(f"Generated presentation file size: {file_size} bytes (images: {self.include_images})")
        
        if file_size == 0:
            raise ValueError("Generated presentation file is empty")
            
        return temp_file
    
    def _create_presentation_from_structured_content(self, structured_content: List[Dict[str, Any]], include_images: bool = False):
        """Create a PowerPoint presentation from structured content."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Create a new presentation
        prs = Presentation()
        
        # Set slide dimensions to widescreen (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        for slide_data in structured_content:
            # Create slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title = slide_data.get('title', 'Slide Title')
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(31, 56, 100)
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add content
            content = slide_data.get('content', [])
            if content:
                content_text = '\n'.join([f"• {item}" if not item.startswith('•') else item for item in content])
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(4.5))
                content_frame = content_box.text_frame
                content_frame.text = content_text
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(18)
                content_para.font.color.rgb = RGBColor(64, 64, 64)
        
        return prs