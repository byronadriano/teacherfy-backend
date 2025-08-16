# resources/generators/slide_processor.py - Smart, Language-Agnostic, Multi-Subject Enhanced Version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import logging
import re
import io
from PIL import Image
from pptx.parts.image import Image as PptxImage
from collections import Counter

logger = logging.getLogger(__name__)

# Enhanced presentation styles
STYLE = {
    'colors': {
        'title': RGBColor(44, 62, 80),      # Dark blue-gray for titles
        'body': RGBColor(44, 62, 80),        # Dark blue-gray for body text
        'accent': RGBColor(41, 128, 185)     # Bright blue for emphasis
    },
    'fonts': {
        'title': 'Calibri',
        'body': 'Calibri'
    },
    'sizes': {
        'title': Pt(40),           # Larger title
        'body': Pt(18),           # Readable body text
        'bullet': Pt(16),         # Slightly smaller bullet points
        'notes': Pt(14)           # Comfortable notes size
    }
}

def clean_text_for_presentation(text):
    """
    Clean text specifically for PowerPoint presentations.
    Remove all markdown and formatting while preserving readability.
    """
    if not text or not isinstance(text, str):
        return ""
    
    # Remove markdown bold/italic formatting
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # **bold** -> bold
    text = re.sub(r'\*([^*]+)\*', r'\1', text)      # *italic* -> italic
    text = re.sub(r'__([^_]+)__', r'\1', text)      # __bold__ -> bold
    text = re.sub(r'_([^_]+)_', r'\1', text)        # _italic_ -> italic
    
    # Remove strikethrough
    text = re.sub(r'~~([^~]+)~~', r'\1', text)      # ~~strike~~ -> strike
    
    # Remove markdown headers but keep the text
    text = re.sub(r'^#{1,6}\s*(.+)$', r'\1', text, flags=re.MULTILINE)
    
    # Remove markdown links but keep the text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)  # [text](url) -> text
    
    # Remove inline code backticks
    text = re.sub(r'`([^`]+)`', r'\1', text)        # `code` -> code
    
    # Remove section dividers and markers
    text = re.sub(r'^---+$', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\*\*Section \d+:', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\*\*Slide \d+:', '', text, flags=re.MULTILINE)
    
    # Clean up standalone asterisks
    text = re.sub(r'^\*+\s*', '', text)             # Remove leading asterisks
    text = re.sub(r'\s*\*+$', '', text)             # Remove trailing asterisks
    
    # Clean up bullet points and numbering (but preserve the content)
    text = re.sub(r'^[-•*]\s*', '', text)           # Remove bullet points
    text = re.sub(r'^\d+\.\s*', '', text)           # Remove numbering
    
    # Clean up multiple spaces and normalize whitespace
    text = ' '.join(text.split())
    
    # Remove any remaining special formatting characters
    text = text.replace('---', '')                   # Remove horizontal rules
    text = text.replace('***', '')                   # Remove emphasis combinations
    
    # Remove content labels that might have been included
    if text.lower().startswith('content:'):
        text = text[8:].strip()
    
    return text.strip()

def clean_content_list_for_presentation(content_list):
    """Clean a list of content items for presentation use."""
    if not content_list or not isinstance(content_list, list):
        return []
    
    cleaned_list = []
    for item in content_list:
        if isinstance(item, str):
            cleaned_item = clean_text_for_presentation(item)
            # Skip empty items and content headers
            if cleaned_item and cleaned_item.lower() not in ['content', 'content:', '---']:
                cleaned_list.append(cleaned_item)
    
    return cleaned_list

def extract_enhanced_search_keywords(structured_content):
    """
    Truly smart keyword extraction that works with ANY language and subject without hardcoding.
    Uses statistical analysis and pattern recognition for 20+ subjects and languages.
    """
    # Combine all text content
    all_text = []
    for slide in structured_content:
        if slide.get('title'):
            all_text.append(slide['title'])
        if slide.get('content'):
            all_text.extend(slide['content'])
    
    if not all_text:
        return 'classroom education learning'
    
    text = ' '.join(all_text).lower()
    
    # Step 1: Extract content features using pattern recognition
    content_features = analyze_content_patterns(text)
    
    # Step 2: Extract meaningful terms using statistical methods
    meaningful_terms = extract_statistical_terms(text)
    
    # Step 3: Build search query using extracted features
    search_query = build_smart_search_query(content_features, meaningful_terms)
    
    logger.info(f"Smart extraction - Features: {content_features}, Query: '{search_query}'")
    return search_query

def analyze_content_patterns(text):
    """Analyze content using universal patterns that work across all languages and subjects."""
    features = set()
    
    # Enhanced Math patterns - more comprehensive decimal detection
    if re.search(r'\d+[\+\-\*\/×÷=]\d+', text):
        features.add('mathematics')
    
    if re.search(r'\d+\/\d+', text):  # Fractions
        features.add('fractions')
        features.add('mathematics')
    
    if re.search(r'\d+\.\d+', text):  # Decimals
        features.add('decimals')
        features.add('mathematics')
    
    # Enhanced decimal and place value detection
    if re.search(r'\b(decimal|decimals|tenths|hundredths|thousandths)\b', text, re.IGNORECASE):
        features.add('decimals')
        features.add('mathematics')
    
    if re.search(r'\b(powers?\s+of\s+10|place\s+value|place\s+values)\b', text, re.IGNORECASE):
        features.add('place_value')
        features.add('mathematics')
    
    if re.search(r'\b\d+%\b', text):  # Percentages
        features.add('statistics')
        features.add('mathematics')
    
    # Science patterns
    if re.search(r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b', text):  # Scientific names
        features.add('science')
    
    if re.search(r'\b\d+°[CF]?\b', text):  # Temperature
        features.add('science')
    
    if re.search(r'\b(CO2|H2O|O2|pH|DNA|RNA|NaCl)\b', text, re.IGNORECASE):  # Chemical formulas
        features.add('chemistry')
        features.add('science')
    
    # Biology patterns
    if re.search(r'\b(cell|cells|organism|species|evolution|photosynthesis)\b', text, re.IGNORECASE):
        features.add('biology')
        features.add('science')
    
    # Physics patterns
    if re.search(r'\b(force|energy|motion|gravity|velocity|acceleration)\b', text, re.IGNORECASE):
        features.add('physics')
        features.add('science')
    
    # Data visualization indicators (universal)
    chart_indicators = [r'chart', r'graph', r'table', r'data', r'survey', r'sample', r'diagram']
    if any(re.search(pattern, text, re.IGNORECASE) for pattern in chart_indicators):
        features.add('data_visualization')
    
    # History patterns
    if re.search(r'\b\d{4}\b', text):  # Years
        features.add('history')
    
    if re.search(r'\b(ancient|medieval|century|empire|civilization|war|revolution)\b', text, re.IGNORECASE):
        features.add('history')
    
    # Geography patterns
    if re.search(r'\b[A-Z][a-z]+\s+(river|mountain|ocean|sea|lake|country|continent)\b', text, re.IGNORECASE):
        features.add('geography')
    
    if re.search(r'\b(climate|weather|population|capital|border|map)\b', text, re.IGNORECASE):
        features.add('geography')
    
    # Language Arts patterns
    sentence_count = len(re.findall(r'[.!?]+', text))
    word_count = len(text.split())
    if word_count > 0 and sentence_count / word_count > 0.08:  # High punctuation density
        features.add('language_arts')
    
    if re.search(r'\b(reading|writing|grammar|vocabulary|literature|poetry|story)\b', text, re.IGNORECASE):
        features.add('language_arts')
    
    # Arts patterns
    if re.search(r'\b(art|painting|drawing|sculpture|color|brush|canvas|creative)\b', text, re.IGNORECASE):
        features.add('visual_arts')
    
    if re.search(r'\b(music|song|rhythm|melody|instrument|note|piano|guitar)\b', text, re.IGNORECASE):
        features.add('music')
    
    # Physical Education patterns
    if re.search(r'\b(sport|exercise|fitness|health|running|jumping|team|game)\b', text, re.IGNORECASE):
        features.add('physical_education')
    
    # Technology patterns
    if re.search(r'\b(computer|software|coding|programming|digital|internet|technology)\b', text, re.IGNORECASE):
        features.add('technology')
    
    # Social Studies patterns
    if re.search(r'\b(government|democracy|election|citizen|community|society|culture)\b', text, re.IGNORECASE):
        features.add('social_studies')
    
    # Economics patterns
    if re.search(r'\b(money|economy|business|trade|market|bank|finance)\b', text, re.IGNORECASE):
        features.add('economics')
    
    # Fun/Entertainment patterns
    if re.search(r'\b(fun|game|play|entertainment|hobby|leisure|enjoyment)\b', text, re.IGNORECASE):
        features.add('fun')
    
    # Holiday patterns
    holiday_patterns = [
        r'\b(christmas|halloween|thanksgiving|easter|valentine|birthday)\b',
        r'\b(holiday|celebration|festival|party|tradition|seasonal)\b',
        r'\b(december|january|october|november|february|march|april)\b'
    ]
    if any(re.search(pattern, text, re.IGNORECASE) for pattern in holiday_patterns):
        features.add('holidays')
    
    # Seasonal patterns
    if re.search(r'\b(spring|summer|fall|autumn|winter|season)\b', text, re.IGNORECASE):
        features.add('seasonal')
    
    # Special subjects
    if re.search(r'\b(special|therapy|intervention|support|accommodation)\b', text, re.IGNORECASE):
        features.add('special_education')
    
    # Health patterns
    if re.search(r'\b(health|nutrition|diet|wellness|safety|hygiene|medical)\b', text, re.IGNORECASE):
        features.add('health')
    
    # Environmental patterns
    if re.search(r'\b(environment|nature|ecology|conservation|sustainability|green)\b', text, re.IGNORECASE):
        features.add('environmental')
    
    return features

def extract_statistical_terms(text):
    """Extract meaningful terms using statistical analysis instead of stop word lists."""
    # Extract all words with Unicode support for any language
    words = re.findall(r'\b[\w\u00C0-\u024F\u1E00-\u1EFF\u0100-\u017F\u0180-\u024F]+\b', text.lower())
    
    # Statistical filtering instead of hardcoded stop words
    word_freq = Counter(words)
    total_words = len(words)
    
    if total_words == 0:
        return []
    
    meaningful_terms = []
    
    for word, count in word_freq.items():
        # Skip if word is too common or too rare
        frequency_ratio = count / total_words
        
        # Filter criteria (language-agnostic)
        if (len(word) >= 4 and                          # Reasonable length
            frequency_ratio < 0.3 and                  # Not too common (like "the", "de", etc.)
            frequency_ratio > (1 / total_words) and    # Not hapax legomena
            not word.isdigit() and                     # Not pure numbers
            not re.match(r'^(st|nd|rd|th|er|ème|º|ª|ый|ая|ое)$', word)):  # Not ordinal suffixes
            meaningful_terms.append(word)
    
    # Return top terms by a combination of frequency and length
    scored_terms = []
    for term in meaningful_terms:
        # Score based on length (longer = more specific) and moderate frequency
        freq_score = word_freq[term]
        length_score = min(len(term), 12)  # Cap at 12
        uniqueness_score = 1 / (word_freq[term] + 1)  # Prefer less common but not unique
        
        total_score = length_score * uniqueness_score * freq_score
        scored_terms.append((term, total_score))
    
    # Sort by score and return top 5
    top_terms = sorted(scored_terms, key=lambda x: x[1], reverse=True)[:5]
    return [term for term, score in top_terms]

def build_smart_search_query(features, terms):
    """Build search query based on detected features and extracted terms for 20+ subjects."""
    
    # Map features to search contexts (comprehensive subject mapping with better math contexts)
    feature_contexts = {
        # Enhanced Math Subjects
        'decimals': 'decimal numbers mathematics place value classroom',
        'place_value': 'place value chart mathematics tens hundreds classroom',
        'fractions': 'fractions mathematics visual circles pie charts',
        'mathematics': 'mathematics classroom education numbers calculations',
        'statistics': 'data charts graphs statistics mathematics',
        'data_visualization': 'charts graphs data visualization classroom',
        
        # Core Academic Subjects
        'science': 'science classroom experiment laboratory education',
        'chemistry': 'chemistry science laboratory beakers molecules',
        'biology': 'biology science nature organisms cells',
        'physics': 'physics science motion energy forces',
        'history': 'history timeline education classroom historical',
        'geography': 'geography maps world classroom globe',
        'language_arts': 'reading books classroom library education',
        
        # Arts and Creative Subjects
        'visual_arts': 'art painting creativity classroom colorful artistic',
        'music': 'music instruments classroom education musical',
        
        # Physical and Health
        'physical_education': 'sports exercise fitness gymnasium education',
        'health': 'health wellness education classroom medical',
        
        # Modern Subjects
        'technology': 'technology computers digital classroom programming',
        'social_studies': 'community society classroom education civic',
        'economics': 'money economics business classroom finance',
        'environmental': 'environment nature conservation classroom green',
        
        # Special Categories
        'fun': 'fun games colorful playful educational activities',
        'holidays': 'holiday celebration festive colorful seasonal',
        'seasonal': 'seasonal nature classroom decorative calendar',
        'special_education': 'inclusive education support classroom diverse',
        
        # Temporal
        'time_based': 'timeline calendar education classroom chronological'
    }
    
    # Start with detected features
    search_parts = []
    
    # Priority order for feature selection (most specific first)
    priority_features = [
        # Very specific math subjects first (for better relevance)
        'decimals', 'place_value', 'fractions', 'statistics', 'data_visualization',
        
        # Specific science subjects
        'chemistry', 'biology', 'physics',
        
        # Creative and special subjects
        'visual_arts', 'music', 'physical_education', 'technology', 
        'economics', 'environmental', 'fun', 'holidays', 'seasonal', 'special_education',
        
        # General subjects last (fallback)
        'mathematics', 'science', 'history', 'geography', 'language_arts',
        'social_studies', 'health', 'time_based'
    ]
    
    # Add primary feature context (first match wins for specificity)
    for feature in priority_features:
        if feature in features:
            search_parts.append(feature_contexts[feature])
            logger.debug(f"Selected primary feature: {feature}")
            break
    
    # Add meaningful extracted terms (max 2 for focused results)
    if terms:
        # Filter terms to avoid generic educational words
        generic_terms = {'lesson', 'class', 'student', 'learn', 'study', 'education', 'school', 'today', 'will'}
        
        filtered_terms = []
        for term in terms[:3]:  # Check top 3 terms
            if term not in generic_terms and len(term) > 4:
                filtered_terms.append(term)
                if len(filtered_terms) >= 2:  # Limit to 2 terms
                    break
        
        if filtered_terms:
            search_parts.extend(filtered_terms)
            logger.debug(f"Added extracted terms: {filtered_terms}")
    
    # Always add educational context if not present
    if not any('classroom' in part or 'education' in part for part in search_parts):
        search_parts.append('education classroom')
    
    # Join and deduplicate
    search_query = ' '.join(search_parts)
    words = search_query.split()
    
    # Remove duplicates while preserving order
    seen = set()
    unique_words = []
    for word in words:
        if word not in seen:
            seen.add(word)
            unique_words.append(word)
    
    result = ' '.join(unique_words)
    
    # Fallback if something went wrong
    if len(result.strip()) < 10:
        result = 'education classroom colorful learning mathematics'
    
    logger.debug(f"Final search query: '{result}'")
    return result

def _enhance_structured_content_for_presentation(structured_content):
    """Enhance structured content for better PowerPoint presentation processing"""
    
    if not structured_content or not isinstance(structured_content, list):
        return structured_content
    
    enhanced_content = []
    
    for slide_data in structured_content:
        if not isinstance(slide_data, dict):
            continue
            
        enhanced_slide = {
            'title': slide_data.get('title', 'Untitled'),
            'layout': slide_data.get('layout', 'TITLE_AND_CONTENT'),
            'content': []
        }
        
        # Process content intelligently
        raw_content = slide_data.get('content', [])
        
        if isinstance(raw_content, list):
            for item in raw_content:
                if isinstance(item, str):
                    # Clean and filter content
                    cleaned_item = clean_text_for_presentation(item)
                    if cleaned_item and not _is_metadata_content(cleaned_item):
                        enhanced_slide['content'].append(cleaned_item)
        
        # Add metadata as separate fields for processing
        if 'structured_questions' in slide_data:
            enhanced_slide['structured_questions'] = slide_data['structured_questions']
        
        if 'teacher_notes' in slide_data:
            enhanced_slide['teacher_notes'] = slide_data['teacher_notes']
            
        if 'differentiation_tips' in slide_data:
            enhanced_slide['differentiation_tips'] = slide_data['differentiation_tips']
        
        enhanced_content.append(enhanced_slide)
    
    logger.info(f"Enhanced {len(structured_content)} slides for presentation processing")
    return enhanced_content

def _is_metadata_content(content_text):
    """Check if content is metadata that shouldn't appear on slides"""
    
    metadata_indicators = [
        'teacher note:',
        'differentiation tip:',
        'assessment check:',
        'for teachers:',
        'instructor guidance:',
        'teaching strategy:',
        'lesson plan:'
    ]
    
    content_lower = content_text.lower().strip()
    
    return any(content_lower.startswith(indicator) for indicator in metadata_indicators)

def add_image_to_slide(slide, image_bytes, lesson_topic=""):
    """
    Add an image to a widescreen slide (13.33" x 7.5") with proper positioning.
    Places image on the right side, centered vertically with text content.
    """
    try:
        # Create a BytesIO object from the image bytes
        image_stream = io.BytesIO(image_bytes)
        
        # Open image with PIL to get dimensions and potentially resize
        with Image.open(image_stream) as img:
            original_width, original_height = img.size
            
          # Get actual slide dimensions
            try:
                if hasattr(slide, 'parent') and slide.parent:
                    slide_width = slide.parent.slide_width
                    slide_height = slide.parent.slide_height
                    logger.debug(f"Using actual slide dimensions: {slide_width} x {slide_height}")
                else:
                    raise AttributeError("No parent presentation available")
            except Exception as e:
                # Fallback to common dimensions
                slide_width = Inches(13.33)
                slide_height = Inches(7.5)
                logger.debug(f"Using fallback dimensions due to: {e}")
            
            # Calculate target image size as percentage of slide
            target_width = slide_width * 0.32   # 32% of slide width
            target_height = slide_height * 0.45  # 45% of slide height
            
            # Calculate aspect ratio and adjust if needed
            img_aspect = original_width / original_height
            target_aspect = target_width / target_height
            
            if img_aspect > target_aspect:
                # Image is wider than target, fit by width
                final_width = target_width
                final_height = target_width / img_aspect
            else:
                # Image is taller than target, fit by height
                final_height = target_height
                final_width = target_height * img_aspect
            
            # Position on right side of slide, centered vertically
            left = slide_width - final_width - (slide_width * 0.04)  # 4% margin from right
            
            # Center vertically in the content area (below title)
            content_start = slide_height * 0.25  # Start at 25% of slide height
            content_height = slide_height * 0.6   # Use 60% of slide height for content
            content_center = content_start + (content_height / 2)
            top = content_center - (final_height / 2)
            
            # Ensure image doesn't go too high or too low
            min_top = slide_height * 0.22  # Don't overlap with title
            max_top = slide_height * 0.85 - final_height  # Don't go off bottom
            
            if top < min_top:
                top = min_top
            elif top > max_top:
                top = max_top
            
            # Reset image stream position
            image_stream.seek(0)
            
            # Add image to slide
            picture = slide.shapes.add_picture(
                image_stream, 
                left, 
                top, 
                final_width, 
                final_height
            )
            
            # Add subtle styling for professional look            
            try:
                line = picture.line
                line.color.rgb = RGBColor(200, 200, 200)  # Light gray border
                line.width = Pt(0.75)
            except:
                pass  # Skip styling if it causes issues
            
            logger.info(f"Successfully added image to slide (size: {final_width} x {final_height}, position: right-center)")
            return True
            
    except Exception as e:
        logger.error(f"Failed to add image to slide: {e}")
        return False
    finally:
        if 'image_stream' in locals():
            image_stream.close()

def add_text_box_to_slide(slide, content_items, with_image=False):
    """Add a text box to slide with dynamic sizing based on actual slide dimensions."""
    
    # Get actual slide dimensions from the presentation
    try:
        if hasattr(slide, 'parent') and slide.parent:
            slide_width = slide.parent.slide_width
            slide_height = slide.parent.slide_height
            logger.debug(f"Actual slide dimensions: {slide_width} x {slide_height}")
        else:
            raise AttributeError("No parent presentation available")
    except Exception as e:
        # Fallback to common widescreen dimensions
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
        logger.debug(f"Using fallback dimensions due to: {e}")
    
    if with_image:
        # Text takes up left portion, leaving right for image
        left = Inches(0.8)
        top = Inches(2.1)
        width = slide_width * 0.55  # 55% of slide width
        height = slide_height * 0.55  # 55% of slide height
    else:
        # Use most of slide when no image
        left = Inches(0.8)
        top = Inches(2.1)
        width = slide_width * 0.85  # 85% of slide width
        height = slide_height * 0.6   # 60% of slide height
    
    # Ensure minimum and maximum dimensions
    width = max(min(width, Inches(11)), Inches(4))
    height = max(min(height, Inches(5)), Inches(3))
    
    # Create text box with dynamic dimensions
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    
    # Optimize text frame properties
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_bottom = Inches(0.15)
    text_frame.word_wrap = True
    text_frame.auto_size = None  # Prevent auto-sizing
    
    # Use cleaned content with appropriate text sizing
    cleaned_items = clean_content_list_for_presentation(content_items)
    
    for item in cleaned_items:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = STYLE['fonts']['body']
        p.font.size = Pt(18)              # Standard readable size
        p.font.color.rgb = STYLE['colors']['body']
        p.space_after = Pt(6)             # Good spacing for readability
        p.line_spacing = 1.2              # Good line spacing
        p.level = 0                       # Consistent bullet level
    
    layout_desc = 'with image accommodation' if with_image else 'full slide'
    logger.info(f"Added text box ({layout_desc}) - size: {width/914400:.1f}\"x{height/914400:.1f}\"")

def clear_all_placeholder_content(slide):
    """Safely clear placeholder content without corrupting the presentation structure."""
    try:
        placeholders_cleared = 0
        
        for shape in slide.shapes:
            try:
                # Skip title shapes to preserve them
                if hasattr(slide.shapes, 'title') and shape == slide.shapes.title:
                    continue
                
                # Only clear text-based placeholders safely
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # Safe method: only clear text content, not structure
                        try:
                            shape.text_frame.clear()
                            placeholders_cleared += 1
                            logger.debug(f"Safely cleared placeholder shape")
                        except:
                            # If clear() fails, try setting empty text
                            try:
                                if hasattr(shape.text_frame, 'text'):
                                    shape.text_frame.text = ""
                                    placeholders_cleared += 1
                            except:
                                pass
                
                # Check by placeholder format (safer approach)
                elif hasattr(shape, 'placeholder_format'):
                    try:
                        placeholder_type = getattr(shape.placeholder_format, 'type', None)
                        if placeholder_type in [2, 7, 8, 14]:  # Content placeholders
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                shape.text_frame.clear()
                                placeholders_cleared += 1
                                logger.debug(f"Cleared placeholder by type: {placeholder_type}")
                    except:
                        pass
                        
            except Exception as e:
                logger.debug(f"Could not process shape safely: {e}")
                continue
        
        logger.info(f"Safely cleared {placeholders_cleared} placeholder shapes")
        return placeholders_cleared > 0
        
    except Exception as e:
        logger.warning(f"Error in safe placeholder clearing: {e}")
        return False

def find_content_placeholder(slide):
    """Find a suitable content placeholder on the slide"""
    # Try to find placeholders in order of preference
    for placeholder in slide.placeholders:
        try:
            # Check if it's a content placeholder (not title)
            if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type:
                placeholder_type = placeholder.placeholder_format.type
                # Look for content, body, or object placeholders
                if placeholder_type in [2, 7, 8, 14]:  # CONTENT, BODY, OBJECT, CONTENT_WITH_CAPTION
                    return placeholder
        except:
            continue
    
    # If no specific content placeholder found, try by index
    try:
        # Most templates have content at index 1
        if len(slide.placeholders) > 1:
            return slide.placeholders[1]
    except:
        pass
    
    # Last resort: find any text-capable placeholder that's not the title
    for i, placeholder in enumerate(slide.placeholders):
        try:
            if i == 0:  # Skip title placeholder
                continue
            # Test if we can access text_frame
            if hasattr(placeholder, 'text_frame'):
                return placeholder
        except:
            continue
    
    return None

def create_clean_presentation_with_images(structured_content, include_images=False):
    """Create a PowerPoint presentation from clean structured content with enhanced images"""
    # Initialize for enhanced per-slide image processing
    
    # Enhanced content processing for JSON structured data
    processed_content = _enhance_structured_content_for_presentation(structured_content)
    
    # Prefer the repo-level static templates location, fall back to module templates
    template_candidates = [
        os.path.join('static', 'templates', 'FINAL_base_template_v1.pptx'),
        os.path.join(os.path.dirname(__file__), 'templates', 'FINAL_base_template_v1.pptx'),
        os.path.join('templates', 'FINAL_base_template_v1.pptx'),
        os.path.join('static', 'templates', 'base_template.pptx'),
        os.path.join(os.path.dirname(__file__), 'templates', 'base_template.pptx'),
        'base_template.pptx'
    ]

    template_path = None
    for candidate in template_candidates:
        if os.path.exists(candidate):
            template_path = candidate
            break

    if not template_path:
        # Last resort: empty Presentation
        template_path = None
    
    # Create presentation
    try:
        prs = Presentation(template_path)
        logger.info(f"Using template: {template_path}")
    except Exception as e:
        logger.warning(f"Could not load template: {e}. Creating blank presentation.")
        prs = Presentation()
    
    # Initialize Unsplash service if images are requested
    unsplash_service = None
    
    if include_images:
        try:
            # Import the global Unsplash service from core.services (project layout)
            from core.services.unsplash_service import unsplash_service as us
            unsplash_service = us
            logger.info("Unsplash service initialized for per-slide image generation")
        except Exception as e:
            logger.error(f"Error initializing Unsplash service: {e}")
            include_images = False  # Disable images for this generation
    
    # Log available slide layouts for debugging
    logger.info(f"Available slide layouts: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        try:
            layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
            logger.debug(f"Layout {i}: {layout_name} - {len(layout.placeholders)} placeholders")
        except:
            logger.debug(f"Layout {i}: Unknown layout")
    
    # Process each slide with clean structure and improved layout
    for slide_index, slide_data in enumerate(processed_content):
        try:
            # Try different layouts in order of preference
            layout_indices_to_try = [1, 0, 2, 3, 4]  # Title+Content, Title, Section, etc.
            slide = None
            
            for layout_idx in layout_indices_to_try:
                if layout_idx < len(prs.slide_layouts):
                    try:
                        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                        logger.debug(f"Successfully used layout {layout_idx} for slide {slide_index + 1}")
                        break
                    except Exception as e:
                        logger.debug(f"Failed to use layout {layout_idx}: {e}")
                        continue
            
            if not slide:
                # Fallback to the first available layout
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                logger.warning(f"Using fallback layout 0 for slide {slide_index + 1}")
            
            # Add contextually relevant image to each content slide
            has_image = False
            if include_images and unsplash_service and slide_index > 0:  # Skip learning objectives slide
                try:
                    # Create optimized search query from slide title and content
                    slide_title = slide_data.get('title', '')
                    slide_content = slide_data.get('content', [])
                    
                    # Use enhanced image search logic
                    search_query = generate_optimized_image_search_query(slide_title, slide_content)
                    
                    if search_query.strip():
                        logger.info(f"Searching for image for slide {slide_index + 1} with query: '{search_query}'")
                        
                        # Search for slide-specific image
                        photo_data = unsplash_service.search_photo(search_query)
                        if photo_data:
                            # Download image
                            image_bytes = unsplash_service.download_photo(photo_data)
                            if image_bytes:
                                has_image = add_image_to_slide(slide, image_bytes, slide_title)
                                
                                # Add attribution to this slide
                                try:
                                    attribution_text = f"Photo: {photo_data['photographer_name']} on Unsplash"
                                    attr_box = slide.shapes.add_textbox(
                                        Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.3)
                                    )
                                    attr_frame = attr_box.text_frame
                                    attr_para = attr_frame.add_paragraph()
                                    attr_para.text = attribution_text
                                    attr_para.font.size = Pt(8)
                                    attr_para.font.color.rgb = RGBColor(128, 128, 128)
                                    attr_para.alignment = PP_ALIGN.RIGHT
                                except Exception as attr_e:
                                    logger.warning(f"Failed to add attribution to slide {slide_index + 1}: {attr_e}")
                                
                                logger.info(f"Added contextual image to slide {slide_index + 1}: '{slide_title}' by {photo_data['photographer_name']}")
                            else:
                                logger.warning(f"Failed to download image for slide {slide_index + 1}")
                        else:
                            logger.warning(f"No image found for slide {slide_index + 1} query: '{search_query}'")
                    
                except Exception as e:
                    logger.error(f"Error adding image to slide {slide_index + 1}: {e}")
                    has_image = False
            
            # Clean and add title
            title_added = False
            raw_title = slide_data.get('title', 'Untitled')
            clean_title = clean_text_for_presentation(raw_title)
            
            if slide.shapes.title:
                try:
                    title_frame = slide.shapes.title.text_frame
                    title_frame.clear()
                    title_para = title_frame.add_paragraph()
                    title_para.text = clean_title
                    title_para.font.name = STYLE['fonts']['title']
                    title_para.font.size = STYLE['sizes']['title']
                    title_para.font.color.rgb = STYLE['colors']['title']
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER
                    title_added = True
                    logger.debug(f"Added clean title to slide {slide_index + 1}: {clean_title}")
                except Exception as e:
                    logger.warning(f"Failed to add title to slide {slide_index + 1}: {e}")
            
            # WIDESCREEN CONTENT HANDLING - Fixed for 13.33" x 7.5" template
            raw_content_items = slide_data.get('content', [])
            clean_content_items = clean_content_list_for_presentation(raw_content_items)
            
            if clean_content_items:
                # ALWAYS clear placeholders and use text boxes for consistency
                logger.info(f"Using text box for slide {slide_index + 1} (image: {has_image})")
                
                # CRITICAL: Clear ALL placeholders to prevent conflicts
                clear_all_placeholder_content(slide)
                
                # Add our custom text box with widescreen sizing
                add_text_box_to_slide(slide, clean_content_items, has_image)
                
                # Handle title if needed
                if not title_added and clean_title:
                    # Position title for widescreen
                    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(10), Inches(1))
                    title_frame = title_box.text_frame
                    title_para = title_frame.add_paragraph()
                    title_para.text = clean_title
                    title_para.font.name = STYLE['fonts']['title']
                    title_para.font.size = STYLE['sizes']['title']
                    title_para.font.color.rgb = STYLE['colors']['title']
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER
                    logger.debug(f"Added widescreen title as text box to slide {slide_index + 1}")
        
        except Exception as e:
            logger.error(f"Error creating slide {slide_index + 1}: {e}")
            # Create a basic slide with text boxes as last resort
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                
                # Add title as text box
                clean_title = clean_text_for_presentation(slide_data.get('title', f'Slide {slide_index + 1}'))
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                title_frame = title_box.text_frame
                title_para = title_frame.add_paragraph()
                title_para.text = clean_title
                title_para.font.name = STYLE['fonts']['title']
                title_para.font.size = STYLE['sizes']['title']
                title_para.font.color.rgb = STYLE['colors']['title']
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                
                # Add content as text box
                content_items = slide_data.get('content', [])
                if content_items:
                    add_text_box_to_slide(slide, content_items, False)
                
                logger.info(f"Created fallback slide {slide_index + 1} with text boxes")
            except Exception as fallback_error:
                logger.error(f"Failed to create fallback slide {slide_index + 1}: {fallback_error}")
    
    # Note: Individual image attributions are now added to each slide that contains an image
    # This provides better context and proper credit for each specific image used
    
    logger.info(f"Created presentation with {len(processed_content)} slides (images: {'enabled' if include_images else 'disabled'})")
    return prs

def generate_optimized_image_search_query(slide_title, slide_content):
    """
    Generate highly optimized search query for a specific slide's image.
    Uses advanced content analysis, subject detection, and visual relevance scoring.
    """
    try:
        # Combine title and content for analysis
        title_text = clean_text_for_presentation(slide_title) if slide_title else ''
        content_text = ' '.join([clean_text_for_presentation(item) for item in slide_content if item])
        combined_text = f"{title_text} {content_text}".lower().strip()
        
        if not combined_text:
            return "classroom education learning"
        
        # Step 1: Analyze content importance - prioritize main content over questions/examples
        main_content, supplementary_content = separate_main_from_supplementary_content(slide_content, title_text)
        
        # Step 2: Extract visual terms with content importance weighting
        visual_terms = extract_smart_visual_terms_weighted(main_content, supplementary_content, combined_text)
        
        # Step 3: Detect subject area for context  
        subject_area = detect_subject_area(combined_text)
        
        # Step 4: Build smart search query using weighted content analysis
        search_query = build_smart_search_query(visual_terms, subject_area, title_text, combined_text)
        
        logger.info(f"Optimized image search for '{title_text}': '{search_query}'")
        return search_query
        
    except Exception as e:
        logger.error(f"Error generating optimized image search: {e}")
        return "classroom education learning"

def separate_main_from_supplementary_content(slide_content, title_text):
    """
    Intelligently separate main educational content from supplementary content like questions.
    Main content should drive image selection, not incidental mentions in questions.
    """
    main_content = []
    supplementary_content = []
    
    for item in slide_content:
        if not item:
            continue
            
        item_lower = item.lower().strip()
        
        # Identify supplementary content patterns (questions, examples at the end)
        is_supplementary = (
            item_lower.startswith(('what is', 'what are', 'how do', 'how does', 'why do', 'why does', 'can you')) or
            item_lower.startswith(('example:', 'for example', 'try this:', 'practice:', 'question:')) or
            '?' in item or
            item_lower.startswith(('discuss', 'think about', 'consider', 'imagine'))
        )
        
        if is_supplementary:
            supplementary_content.append(item)
        else:
            main_content.append(item)
    
    # If we don't have main content, treat everything as main content
    if not main_content:
        main_content = slide_content
        supplementary_content = []
    
    return main_content, supplementary_content

def extract_smart_visual_terms_weighted(main_content, supplementary_content, full_text):
    """
    Extract visual terms with heavy weighting toward main content.
    Only use supplementary content if it reinforces the main theme.
    """
    # Extract terms from main content (high priority)
    main_content_text = ' '.join(main_content).lower()
    main_visual_terms = extract_smart_visual_terms(main_content_text)
    
    # Extract terms from supplementary content (low priority)
    supplementary_content_text = ' '.join(supplementary_content).lower()
    supplementary_visual_terms = extract_smart_visual_terms(supplementary_content_text)
    
    # Determine if supplementary terms reinforce main theme
    reinforcing_terms = []
    if main_visual_terms and supplementary_visual_terms:
        # Check if supplementary terms are related to main terms
        for supp_term in supplementary_visual_terms:
            # Only include supplementary terms that relate to main content theme
            if any(are_terms_related(supp_term, main_term) for main_term in main_visual_terms):
                reinforcing_terms.append(supp_term)
    
    # Detect if main content has only weak terms
    weak_terms = {'author', 'book', 'books', 'story', 'paragraph', 'example', 'details', 'idea', 'point'}
    main_has_strong_terms = any(term not in weak_terms for term in main_visual_terms)
    supplementary_has_strong_terms = any(term not in weak_terms for term in supplementary_visual_terms)
    
    # Smart prioritization: use supplementary if main is weak but supplementary is strong
    if not main_has_strong_terms and supplementary_has_strong_terms:
        logger.info(f"Using supplementary content (main content weak): {supplementary_visual_terms}")
        final_terms = supplementary_visual_terms.copy()
        
        # Also try to extract subject matter from examples
        example_subject_terms = extract_subject_from_examples(supplementary_content)
        for term in example_subject_terms:
            if term not in final_terms:
                final_terms.append(term)
    else:
        # Normal prioritization: main content first
        final_terms = main_visual_terms.copy()
        
        # Add reinforcing terms that aren't already present
        for term in reinforcing_terms:
            if term not in final_terms:
                final_terms.append(term)
    
    # If still no strong terms found, try example subject extraction
    if not final_terms or not any(term not in weak_terms for term in final_terms):
        example_subject_terms = extract_subject_from_examples(supplementary_content)
        if example_subject_terms:
            final_terms = example_subject_terms
        else:
            # Final fallback to title-based analysis
            title_terms = extract_smart_visual_terms(full_text)
            final_terms = title_terms
    
    return final_terms

def are_terms_related(term1, term2):
    """
    Check if two terms are thematically related for content analysis.
    """
    # Define thematic relationships
    related_groups = [
        # Math/food combinations
        {'pizza', 'fractions', 'slices', 'pie', 'cake', 'food'},
        {'coins', 'money', 'dollars', 'cents', 'math'},
        {'clock', 'time', 'hours', 'minutes'},
        
        # Science themes
        {'animals', 'ocean', 'fish', 'marine', 'water'},
        {'plants', 'flowers', 'trees', 'garden', 'nature'},
        {'rainforest', 'ecosystem', 'species', 'oxygen', 'trees', 'animals'},
        {'space', 'planets', 'sun', 'moon', 'stars', 'solar'},
        
        # Social studies themes
        {'maps', 'countries', 'globe', 'geography', 'earth'},
        {'community', 'buildings', 'cities', 'neighborhoods'},
        
        # Arts themes
        {'music', 'instruments', 'piano', 'guitar', 'songs'},
        {'art', 'painting', 'colors', 'brushes', 'canvas'},
        
        # PE themes
        {'sports', 'exercise', 'running', 'swimming', 'dancing'},
        {'soccer', 'basketball', 'football', 'ball', 'games'}
    ]
    
    # Check if both terms belong to the same thematic group
    for group in related_groups:
        if term1 in group and term2 in group:
            return True
    
    return False

def extract_subject_from_examples(supplementary_content):
    """
    Extract the actual subject matter from examples when main content lacks visual terms.
    E.g., "about rainforests" -> extract rainforest-related terms
    """
    subject_terms = []
    
    for item in supplementary_content:
        if not item:
            continue
            
        item_lower = item.lower()
        
        # Look for "about X" or "paragraph about X" patterns
        about_patterns = [
            r'about (\w+)',
            r'paragraph about (\w+)', 
            r'story about (\w+)',
            r'book about (\w+)'
        ]
        
        for pattern in about_patterns:
            matches = re.findall(pattern, item_lower)
            for match in matches:
                # Convert the subject to related visual terms
                subject_visual_terms = get_visual_terms_for_subject(match)
                subject_terms.extend(subject_visual_terms)
        
        # Also extract any strong visual terms directly from examples
        example_visual_terms = extract_smart_visual_terms(item_lower)
        
        # Filter to keep only thematically strong terms (avoid weak terms like "book", "author")
        strong_terms = []
        weak_terms = {'book', 'books', 'author', 'movie', 'story', 'paragraph', 'example', 'details'}
        
        for term in example_visual_terms:
            if term not in weak_terms:
                strong_terms.append(term)
        
        subject_terms.extend(strong_terms)
    
    return list(set(subject_terms))  # Remove duplicates

def get_visual_terms_for_subject(subject):
    """
    Convert a subject name to related visual terms for image search.
    """
    subject_mappings = {
        'rainforest': ['rainforest', 'forest', 'trees', 'jungle', 'nature'],
        'rainforests': ['rainforest', 'forest', 'trees', 'jungle', 'nature'],
        'ocean': ['ocean', 'water', 'waves', 'marine', 'fish'],
        'oceans': ['ocean', 'water', 'waves', 'marine', 'fish'],
        'space': ['space', 'planets', 'stars', 'galaxy', 'astronomy'],
        'animals': ['animals', 'wildlife', 'nature', 'creatures'],
        'plants': ['plants', 'flowers', 'garden', 'nature', 'green'],
        'weather': ['weather', 'clouds', 'rain', 'storm', 'sky'],
        'sports': ['sports', 'exercise', 'athletes', 'competition'],
        'music': ['music', 'instruments', 'musical', 'sound'],
        'art': ['art', 'painting', 'creative', 'colors'],
        'science': ['science', 'laboratory', 'research', 'discovery'],
        'math': ['mathematics', 'numbers', 'calculation', 'problem'],
        'mathematics': ['mathematics', 'numbers', 'calculation', 'problem'],
        'history': ['history', 'historical', 'ancient', 'past'],
        'geography': ['geography', 'maps', 'world', 'countries']
    }
    
    return subject_mappings.get(subject, [subject])

def extract_smart_visual_terms(text):
    """
    Smart extraction of actual visual terms mentioned in content.
    Focuses on concrete nouns and objects that can be photographed.
    """
    # Extract concrete visual nouns across ALL educational subjects
    visual_nouns = []
    
    # SCIENCE - Animals, nature, body, experiments, space, ecosystems
    science_terms = r'\b(animals?|birds?|fish|bears?|lions?|tigers?|elephants?|dogs?|cats?|insects?|butterflies?|frogs?|plants?|trees?|flowers?|leaves?|roots?|seeds?|sun|moon|stars?|planets?|earth|rocks?|minerals?|water|ocean|river|lake|mountains?|volcanoes?|clouds?|rain|snow|weather|skeleton|heart|lungs|muscles?|eyes?|microscope|telescope|beaker|laboratory|magnet|battery|rainforest|rainforests?|ecosystem|ecosystems?|species|oxygen|forest|forests?|jungle|nature|wildlife|habitat|habitats?)\b'
    
    # SOCIAL STUDIES - Geography, community, history, culture  
    social_studies_terms = r'\b(maps?|globe|countries?|continents?|cities?|towns?|neighborhoods?|houses?|buildings?|schools?|libraries?|hospitals?|stores?|farms?|factories?|bridges?|roads?|flags?|monuments?|castles?|pyramids?|museums?|communities?|families?)\b'
    
    # PHYSICAL EDUCATION - Sports, exercise, health
    pe_terms = r'\b(balls?|soccer|basketball|football|baseball|tennis|volleyball|swimming|running|jumping|dancing|gymnastics|bikes?|bicycles?|sports?|games?|playground|gym|exercise|yoga|stretching)\b'
    
    # ARTS - Visual arts, music, performance
    arts_terms = r'\b(painting|drawings?|colors?|brushes?|canvas|crayons?|markers?|art|sculptures?|pottery|music|instruments?|piano|guitar|violin|drums?|singing|dancing|theater|stage|costumes?)\b'
    
    # LANGUAGE ARTS - Reading, writing, communication
    language_terms = r'\b(books?|stories?|poems?|newspapers?|magazines?|letters?|pencils?|pens?|paper|notebooks?|writing|reading|library|authors?|characters?)\b'
    
    # MATH - Numbers, shapes, tools (keep math terms but expand)
    math_terms = r'\b(numbers?|shapes?|circles?|triangles?|squares?|rectangles?|cubes?|spheres?|rulers?|calculators?|clocks?|time|money|coins?|dollars?|pizza|pie|cake|fractions?|graphs?|charts?)\b'
    
    # EVERYDAY OBJECTS - Common items that appear in any subject  
    everyday_terms = r'\b(phones?|doors?|windows?|cars?|trucks?|trains?|airplanes?|boats?|computers?|tablets?|cameras?|toys?|tools?|machines?|wheels?)\b'
    
    # FOOD - Very visual and engaging across subjects
    food_terms = r'\b(food|fruits?|vegetables?|apples?|oranges?|bananas?|carrots?|broccoli|bread|sandwiches?|milk|water|juice|pizza|hamburgers?|salad|soup)\b'
    
    # Extract all matches from all subjects
    all_patterns = [
        science_terms, social_studies_terms, pe_terms, arts_terms,
        language_terms, math_terms, everyday_terms, food_terms
    ]
    
    for pattern in all_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        visual_nouns.extend([match.lower() for match in matches])
    
    return list(set(visual_nouns))  # Remove duplicates

def detect_subject_area(text):
    """
    Detect the primary subject area for educational context.
    """
    subject_indicators = {
        'mathematics': ['math', 'number', 'add', 'subtract', 'multiply', 'divide', 'equation', 'solve', 'calculate'],
        'science': ['science', 'experiment', 'observe', 'hypothesis', 'data', 'research', 'discovery'],
        'reading': ['read', 'story', 'book', 'character', 'plot', 'author', 'literature', 'poem'],
        'writing': ['write', 'essay', 'paragraph', 'sentence', 'grammar', 'spelling', 'vocabulary'],
        'social_studies': ['community', 'history', 'geography', 'culture', 'government', 'citizen'],
        'art': ['create', 'draw', 'paint', 'design', 'artistic', 'creative', 'imagination'],
        'physical_education': ['exercise', 'movement', 'sport', 'healthy', 'fitness', 'active'],
        'health': ['healthy', 'nutrition', 'safety', 'hygiene', 'wellness', 'medical']
    }
    
    subject_scores = {}
    for subject, indicators in subject_indicators.items():
        score = sum(1 for indicator in indicators if re.search(r'\b' + re.escape(indicator) + r'\b', text, re.IGNORECASE))
        if score > 0:
            subject_scores[subject] = score
    
    if subject_scores:
        return max(subject_scores, key=subject_scores.get)
    return 'general'

def extract_educational_terms(text):
    """
    Extract specific educational terms that provide context for image selection.
    """
    # Educational grade level indicators
    grade_terms = re.findall(r'\b(?:elementary|primary|kindergarten|grade|student|children|kids)\b', text, re.IGNORECASE)
    
    # Activity type indicators
    activity_terms = re.findall(r'\b(?:hands-on|interactive|group|individual|practice|activity|exercise|game)\b', text, re.IGNORECASE)
    
    # Learning objective terms
    learning_terms = re.findall(r'\b(?:learn|understand|identify|compare|analyze|create|explain|demonstrate)\b', text, re.IGNORECASE)
    
    return {
        'grade_level': list(set(grade_terms)),
        'activity_type': list(set(activity_terms)),
        'learning_objectives': list(set(learning_terms))
    }

def build_smart_search_query(visual_terms, subject_area, title_text, content_text):
    """
    Build search query using actual content terms for maximum relevance.
    Prioritizes concrete visual objects mentioned in the content.
    """
    query_parts = []
    
    # Priority 1: Use actual visual terms from content (most important!)
    if visual_terms:
        # Special combinations that work well together across all subjects
        special_combinations = {
            # MATH combinations
            ('pizza', 'fractions'): 'pizza slices fractions pie chart',
            ('pizza', 'slice'): 'pizza slices fractions pie chart', 
            ('pizza', 'slices'): 'pizza slices fractions pie chart',
            ('sandwich', 'fractions'): 'sandwich half fractions food',
            ('cake', 'fractions'): 'cake slices fractions pie chart',
            ('coin', 'money'): 'coins money counting mathematics',
            ('coins', 'money'): 'coins money counting mathematics',
            ('clock', 'time'): 'clock time telling education',
            
            # SCIENCE combinations
            ('animals', 'ocean'): 'ocean animals fish marine life',
            ('fish', 'ocean'): 'ocean fish marine life underwater',
            ('birds', 'trees'): 'birds trees forest nature wildlife',
            ('plants', 'flowers'): 'flowers plants garden nature blooming',
            ('sun', 'weather'): 'sun weather sky sunshine nature',
            ('clouds', 'rain'): 'clouds rain weather sky storm',
            ('microscope', 'laboratory'): 'microscope laboratory science research',
            ('skeleton', 'bones'): 'skeleton bones human body anatomy',
            ('planets', 'space'): 'planets space solar system astronomy',
            
            # SOCIAL STUDIES combinations
            ('maps', 'countries'): 'world maps countries geography atlas',
            ('globe', 'earth'): 'globe earth world geography planet',
            ('flags', 'countries'): 'flags countries nations patriotic symbols',
            ('buildings', 'cities'): 'buildings cities urban architecture skyline',
            ('farms', 'food'): 'farms agriculture food crops farming',
            ('monuments', 'history'): 'monuments historical landmarks architecture',
            
            # PHYSICAL EDUCATION combinations  
            ('soccer', 'ball'): 'soccer ball sports field football',
            ('basketball', 'gym'): 'basketball gym sports court indoor',
            ('swimming', 'water'): 'swimming pool water sports exercise',
            ('running', 'exercise'): 'running exercise fitness outdoor sports',
            ('dancing', 'music'): 'dancing music performance arts movement',
            
            # ARTS combinations
            ('painting', 'colors'): 'painting colors art brushes canvas colorful',
            ('music', 'instruments'): 'musical instruments music orchestra band',
            ('piano', 'music'): 'piano music keyboard musical instrument',
            ('guitar', 'music'): 'guitar music strings musical instrument',
            ('art', 'brushes'): 'art brushes painting creative supplies',
            
            # LANGUAGE ARTS combinations
            ('books', 'reading'): 'books reading library literature education',
            ('stories', 'books'): 'storybooks reading children literature',
            ('writing', 'pencils'): 'writing pencils paper education school',
            ('library', 'books'): 'library books reading shelves education'
        }
        
        # Check for special combinations first
        combination_found = False
        for combo, query in special_combinations.items():
            if all(term in visual_terms for term in combo):
                query_parts.append(query)
                combination_found = True
                break
        
        # If no special combination, use the most relevant visual terms
        if not combination_found:
            # Check if title gives us context for prioritization
            title_lower = title_text.lower() if title_text else ''
            
            # Context-based priorities
            geometry_context = any(word in title_lower for word in ['geometry', 'shape', 'shapes'])
            time_context = any(word in title_lower for word in ['time', 'telling', 'clock'])
            
            priority_terms = []
            
            # Prioritize based on title context
            if geometry_context:
                geometry_terms = ['phone', 'door', 'notebook', 'book', 'wheel', 'clock']
                for term in visual_terms:
                    if term in geometry_terms:
                        priority_terms.append(term)
                        if len(priority_terms) >= 2:
                            break
                # Add geometric context
                if priority_terms:
                    query_parts.append(' '.join(priority_terms) + ' geometric shapes')
                    combination_found = True
                    
            elif time_context:
                time_terms = ['clock']
                for term in visual_terms:
                    if term in time_terms:
                        priority_terms.append(term)
                        break
                if priority_terms:
                    query_parts.append(' '.join(priority_terms) + ' time telling')
                    combination_found = True
            
            # Subject-aware prioritization based on detected terms
            if not combination_found:
                # Categorize visual terms by subject for smart prioritization
                science_terms = ['animals', 'birds', 'fish', 'plants', 'trees', 'flowers', 'sun', 'moon', 'stars', 'planets', 'ocean', 'mountains', 'microscope', 'telescope']
                pe_terms = ['soccer', 'basketball', 'football', 'tennis', 'swimming', 'running', 'dancing', 'exercise', 'sports', 'gym']
                arts_terms = ['painting', 'music', 'piano', 'guitar', 'violin', 'art', 'colors', 'brushes', 'canvas', 'instruments']
                social_terms = ['maps', 'globe', 'countries', 'cities', 'buildings', 'farms', 'flags', 'monuments', 'communities']
                math_terms = ['numbers', 'shapes', 'circles', 'triangles', 'calculator', 'money', 'coins', 'clocks', 'charts']
                language_terms = ['books', 'stories', 'reading', 'writing', 'pencils', 'library', 'newspapers']
                food_terms = ['pizza', 'cake', 'sandwich', 'apple', 'fruits', 'vegetables', 'food']
                
                # Find the dominant subject based on visual terms present
                subject_scores = {
                    'science': sum(1 for term in visual_terms if term in science_terms),
                    'pe': sum(1 for term in visual_terms if term in pe_terms),  
                    'arts': sum(1 for term in visual_terms if term in arts_terms),
                    'social': sum(1 for term in visual_terms if term in social_terms),
                    'math': sum(1 for term in visual_terms if term in math_terms),
                    'language': sum(1 for term in visual_terms if term in language_terms),
                    'food': sum(1 for term in visual_terms if term in food_terms)
                }
                
                dominant_subject = max(subject_scores, key=subject_scores.get) if max(subject_scores.values()) > 0 else 'general'
                
                # Prioritize terms based on dominant subject
                if dominant_subject == 'science':
                    priority_order = [science_terms, food_terms, math_terms]
                elif dominant_subject == 'pe':
                    priority_order = [pe_terms, food_terms, science_terms]
                elif dominant_subject == 'arts':
                    priority_order = [arts_terms, food_terms, science_terms] 
                elif dominant_subject == 'social':
                    priority_order = [social_terms, food_terms, science_terms]
                elif dominant_subject == 'math':
                    priority_order = [math_terms, food_terms, science_terms]
                elif dominant_subject == 'language':
                    priority_order = [language_terms, food_terms, science_terms]
                else:
                    # General prioritization - food first (most visual), then others
                    priority_order = [food_terms, science_terms, arts_terms, pe_terms]
                
                # Select terms based on priority order
                for term_group in priority_order:
                    for term in visual_terms:
                        if term in term_group and term not in priority_terms:
                            priority_terms.append(term)
                            if len(priority_terms) >= 2:  # Limit to 2 main terms
                                break
                    if len(priority_terms) >= 2:
                        break
                
                # Add any remaining visual terms if we don't have enough
                for term in visual_terms:
                    if term not in priority_terms:
                        priority_terms.append(term)
                        if len(priority_terms) >= 3:
                            break
                
                if priority_terms:
                    query_parts.append(' '.join(priority_terms))
    
    # Priority 2: Add subject context if no specific visual terms
    if not query_parts:
        subject_contexts = {
            'mathematics': 'mathematics classroom education',
            'science': 'science nature education',
            'reading': 'books reading education', 
            'art': 'art creativity education',
            'general': 'classroom education learning'
        }
        context = subject_contexts.get(subject_area, subject_contexts['general'])
        query_parts.append(context)
    
    # Priority 3: Add educational context
    if not any('education' in part for part in query_parts):
        query_parts.append('education')
    
    # Build final query
    final_query = ' '.join(query_parts)
    
    # Clean up and deduplicate
    words = final_query.split()
    unique_words = []
    seen = set()
    for word in words:
        if word not in seen:
            seen.add(word)
            unique_words.append(word)
    
    result = ' '.join(unique_words[:5])  # Limit to 5 most relevant words
    
    # Ensure minimum quality
    if len(result.strip()) < 5:
        result = 'classroom education learning'
    
    return result

def generate_search_query_from_content(structured_content, fallback="elementary classroom educational"):
    """
    Generate an enhanced search query for Unsplash based on lesson content.
    Uses the new enhanced keyword extraction for better image relevance.
    """
    try:
        if not structured_content or len(structured_content) == 0:
            return fallback
        
        # Use the enhanced keyword extraction
        enhanced_query = extract_enhanced_search_keywords(structured_content)
        
        if enhanced_query and enhanced_query != fallback:
            logger.info(f"Using enhanced search query: '{enhanced_query}'")
            return enhanced_query
        
        # Fallback to original logic if enhanced extraction doesn't find anything
        # Find the first slide with meaningful content (skip title-only slides)
        content_slide = None
        for slide in structured_content:
            content_items = slide.get('content', [])
            if content_items and len(content_items) > 0:
                # Check if this slide has real content, not just titles
                has_real_content = any(
                    len(item.strip()) > 10 and 
                    not item.lower().startswith(('students will', 'today we will', 'objectives'))
                    for item in content_items
                )
                if has_real_content:
                    content_slide = slide
                    break
        
        # If no content slide found, use the first slide
        if not content_slide:
            content_slide = structured_content[0]
        
        # Get title and content from the selected slide
        title = content_slide.get('title', '').lower()
        content_items = content_slide.get('content', [])
        content_text = ' '.join(content_items).lower() if content_items else ''
        
        # Combine title and content for analysis
        combined_text = f"{title} {content_text}".strip()
        
        # Use our smart analysis
        features = analyze_content_patterns(combined_text)
        terms = extract_statistical_terms(combined_text)
        result = build_smart_search_query(features, terms)
        
        if result and result != fallback:
            return result
        
        # Final enhanced fallback
        return "elementary classroom learning educational colorful"
        
    except Exception as e:
        logger.error(f"Error generating search query: {e}")
        return fallback

# BACKWARD COMPATIBILITY - Keep the old function for existing code
def create_presentation(structured_content, resource_type="PRESENTATION"):
    """Legacy function for backward compatibility"""
    # Convert old structure to new clean structure if needed
    clean_content = []
    
    for slide_data in structured_content:
        # Clean all the text content
        clean_slide = {
            'title': clean_text_for_presentation(slide_data.get('title', 'Untitled')),
            'layout': slide_data.get('layout', 'TITLE_AND_CONTENT'),
            'content': clean_content_list_for_presentation(slide_data.get('content', []))
        }
        
        # Handle old structure fields - convert to content and clean
        if not clean_slide['content']:
            # Try to extract from old fields
            if slide_data.get('left_column') or slide_data.get('right_column'):
                combined_content = (slide_data.get('left_column', []) + 
                                 slide_data.get('right_column', []))
                clean_slide['content'] = clean_content_list_for_presentation(combined_content)
            elif slide_data.get('teacher_notes'):
                clean_slide['content'] = clean_content_list_for_presentation(slide_data.get('teacher_notes', []))
        
        clean_content.append(clean_slide)
    
    return create_clean_presentation_with_images(clean_content, include_images=False)

def create_clean_presentation(structured_content):
    """Create a PowerPoint presentation from clean structured content without images"""
    return create_clean_presentation_with_images(structured_content, include_images=False)

def parse_outline_to_structured_content(outline_text, resource_type="PRESENTATION"):
    """Parse outline text into clean structured content"""
    logger.info(f"Parsing outline for resource type: {resource_type}")
    
    # Determine section/slide pattern based on resource type
    if resource_type.upper() == "PRESENTATION":
        section_pattern = r"Slide (\d+):\s*(.*)"
        section_word = "Slide"
    else:
        section_pattern = r"Section (\d+):\s*(.*)"
        section_word = "Section"
    
    # Split by section headers
    sections = []
    current_section = None
    
    lines = outline_text.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check if this is a section/slide header
        match = re.match(section_pattern, line)
        if match:
            # Save previous section
            if current_section:
                sections.append(current_section)
            
            # Start new section
            section_num = match.group(1)
            section_title = match.group(2).strip()
            current_section = {
                "title": clean_text_for_presentation(section_title),
                "layout": "TITLE_AND_CONTENT",
                "content": []
            }
        elif line.lower() == "content:":
            # Skip content headers
            continue
        elif line.startswith('-') or line.startswith('•'):
            # This is content
            if current_section:
                clean_content = clean_text_for_presentation(line.lstrip('-•').strip())
                if clean_content:
                    current_section["content"].append(clean_content)
        elif current_section and line:
            # Any other non-empty line goes to content
            clean_content = clean_text_for_presentation(line)
            if clean_content:
                current_section["content"].append(clean_content)
    
    # Don't forget the last section
    if current_section:
        sections.append(current_section)
    
    # If no sections found, create a fallback
    if not sections:
        sections.append({
            "title": "Generated Content",
            "layout": "TITLE_AND_CONTENT",
            "content": [clean_text_for_presentation(line.strip()) for line in lines if line.strip()]
        })
    
    logger.info(f"Successfully parsed {len(sections)} sections for {resource_type}")
    return sections