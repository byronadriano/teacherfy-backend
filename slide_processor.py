from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import os
import logging

logger = logging.getLogger(__name__)

# Enhanced presentation styles with Comic Sans MS and larger title
STYLE = {
    'colors': {
        'title': RGBColor(44, 62, 80),      # Dark blue-gray for titles
        'body': RGBColor(44, 62, 80),        # Dark blue-gray for body text
        'accent': RGBColor(41, 128, 185)     # Bright blue for emphasis
    },
    'fonts': {
        'title': 'Comic Sans MS',
        'body': 'Comic Sans MS'
    },
    'sizes': {
        'title': Pt(50),           # Larger, more prominent title
        'body': Pt(22),           # Readable body text
        'bullet': Pt(20),         # Slightly smaller bullet points
        'notes': Pt(16)           # Comfortable notes size
    }
}

def clean_markdown_formatting(text):
    """Remove markdown formatting while preserving the text structure"""
    # Remove bold markdown indicators
    text = text.replace('**', '')
    
    # Remove any remaining asterisks at the start of lines
    text = text.lstrip('*')
    
    # Clean up any double spaces that might have been created
    text = ' '.join(text.split())
    
    # Handle numbered lists (e.g., "1.", "2.", etc.)
    import re
    text = re.sub(r'^\d+\.\s*', '', text)
    
    return text.strip()

def format_paragraph(paragraph, is_title=False, level=0):
    """Apply consistent formatting to a paragraph"""
    # Clean the text before setting it
    paragraph.text = clean_markdown_formatting(paragraph.text)
    
    paragraph.font.name = STYLE['fonts']['title' if is_title else 'body']
    paragraph.font.size = STYLE['sizes']['title' if is_title else ('body' if level == 0 else 'bullet')]
    paragraph.font.color.rgb = STYLE['colors']['title' if is_title else 'body']
    paragraph.font.bold = is_title or level == 0
    
    if is_title:
        paragraph.alignment = PP_ALIGN.CENTER
    else:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(12 if level > 0 else 20)
        paragraph.space_after = Pt(6)

def parse_outline_to_structured_content(outline_text):
    """Parse the outline text into structured slide content with notes"""
    slides = []
    current_slide = None
    current_section = None
    
    lines = outline_text.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('Slide '):
            if current_slide:
                slides.append(current_slide)
            
            title = line.split(':', 1)[1].strip() if ':' in line else line
            
            layout = "TWO_COLUMN" if any(word in title.lower() for word in 
                ["vs", "comparison", "contrast", "together", "practice"]) else "TITLE_AND_CONTENT"
            
            current_slide = {
                'title': title,
                'layout': layout,
                'content': [],
                'teacher_notes': [],
                'visual_elements': [],
                'left_column': [],
                'right_column': []
            }
            current_section = None
            
        elif line.lower().startswith('content:'):
            current_section = 'content'
        elif line.lower().startswith('teacher notes:'):
            current_section = 'teacher_notes'
        elif line.lower().startswith('visual elements:'):
            current_section = 'visual_elements'
        elif current_section and current_slide:
            if line.startswith(('*', '•', '-')) or line.strip().startswith('•'):
                content = line.lstrip('*•- ').strip()
                if content:
                    if current_slide['layout'] == "TWO_COLUMN" and current_section == 'content':
                        if len(current_slide['left_column']) <= len(current_slide['right_column']):
                            current_slide['left_column'].append(content)
                        else:
                            current_slide['right_column'].append(content)
                    else:
                        current_slide[current_section].append(content)
            elif not line.lower().endswith(':'):
                current_slide[current_section].append(line)
    
    if current_slide:
        slides.append(current_slide)
    
    # Post-processing
    for slide in slides:
        if slide['layout'] == "TWO_COLUMN" and not (slide['left_column'] or slide['right_column']):
            content_length = len(slide['content'])
            mid_point = content_length // 2
            slide['left_column'] = slide['content'][:mid_point]
            slide['right_column'] = slide['content'][mid_point:]
            slide['content'] = []
    
    return slides

def create_presentation(outline_json):
    """Create a PowerPoint presentation with enhanced formatting"""
    template_path = os.path.join(os.path.dirname(__file__), 'templates', 'base_template_finalv2.pptx')
    prs = Presentation(template_path)
    
    for slide_data in outline_json:
        layout_idx = 1 if slide_data['layout'] == "TITLE_AND_CONTENT" else 3
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Format title
        if slide.shapes.title:
            title_frame = slide.shapes.title.text_frame
            title_frame.clear()
            title_para = title_frame.add_paragraph()
            title_para.text = clean_markdown_formatting(slide_data['title'])
            format_paragraph(title_para, is_title=True)
        
        if slide_data['layout'] == "TWO_COLUMN":
            shapes = [shape for shape in slide.placeholders]
            if len(shapes) >= 3:
                left = shapes[1]
                right = shapes[2]
                
                # Left column
                text_frame = left.text_frame
                text_frame.clear()
                for item in slide_data['left_column']:
                    p = text_frame.add_paragraph()
                    p.text = clean_markdown_formatting(item.lstrip('-•* ').strip())
                    format_paragraph(p, level=1)
                
                # Right column
                text_frame = right.text_frame
                text_frame.clear()
                for item in slide_data['right_column']:
                    p = text_frame.add_paragraph()
                    p.text = clean_markdown_formatting(item.lstrip('-•* ').strip())
                    format_paragraph(p, level=1)
        else:
            shapes = [shape for shape in slide.placeholders]
            if len(shapes) >= 2:
                content_placeholder = shapes[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for item in slide_data['content']:
                    p = text_frame.add_paragraph()
                    is_bullet = item.lstrip().startswith(('-', '•', '*'))
                    p.text = clean_markdown_formatting(item.lstrip('-•* ').strip() if is_bullet else item)
                    format_paragraph(p, level=1 if is_bullet else 0)
        
        # Add notes
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame
        notes_text.clear()
        
        if slide_data['teacher_notes']:
            notes_text.text = "Teacher Notes:\n"
            for note in slide_data['teacher_notes']:
                p = notes_text.add_paragraph()
                p.text = f"• {clean_markdown_formatting(note)}"
                p.font.size = STYLE['sizes']['notes']
        
        if slide_data['visual_elements']:
            if notes_text.text:
                p = notes_text.add_paragraph()
                p.text = "\nVisual Elements:"
            else:
                notes_text.text = "Visual Elements:\n"
            
            for visual in slide_data['visual_elements']:
                p = notes_text.add_paragraph()
                p.text = f"• {clean_markdown_formatting(visual)}"
                p.font.size = STYLE['sizes']['notes']
    
    return prs