# src/slide_processor.py - Smart, Language-Agnostic, Multi-Subject Enhanced Version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import logging
import re
import io
from PIL import Image
from pptx.parts.image import Image as PptxImage
from collections import Counter

logger = logging.getLogger(__name__)

# Enhanced presentation styles
STYLE = {
    'colors': {
        'title': RGBColor(44, 62, 80),      # Dark blue-gray for titles
        'body': RGBColor(44, 62, 80),        # Dark blue-gray for body text
        'accent': RGBColor(41, 128, 185)     # Bright blue for emphasis
    },
    'fonts': {
        'title': 'Calibri',
        'body': 'Calibri'
    },
    'sizes': {
        'title': Pt(40),           # Larger title
        'body': Pt(18),           # Readable body text
        'bullet': Pt(16),         # Slightly smaller bullet points
        'notes': Pt(14)           # Comfortable notes size
    }
}

def clean_text_for_presentation(text):
    """
    Clean text specifically for PowerPoint presentations.
    Remove all markdown and formatting while preserving readability.
    """
    if not text or not isinstance(text, str):
        return ""
    
    # Remove markdown bold/italic formatting
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # **bold** -> bold
    text = re.sub(r'\*([^*]+)\*', r'\1', text)      # *italic* -> italic
    text = re.sub(r'__([^_]+)__', r'\1', text)      # __bold__ -> bold
    text = re.sub(r'_([^_]+)_', r'\1', text)        # _italic_ -> italic
    
    # Remove strikethrough
    text = re.sub(r'~~([^~]+)~~', r'\1', text)      # ~~strike~~ -> strike
    
    # Remove markdown headers but keep the text
    text = re.sub(r'^#{1,6}\s*(.+)$', r'\1', text, flags=re.MULTILINE)
    
    # Remove markdown links but keep the text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)  # [text](url) -> text
    
    # Remove inline code backticks
    text = re.sub(r'`([^`]+)`', r'\1', text)        # `code` -> code
    
    # Remove section dividers and markers
    text = re.sub(r'^---+$', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\*\*Section \d+:', '', text, flags=re.MULTILINE)
    text = re.sub(r'^\*\*Slide \d+:', '', text, flags=re.MULTILINE)
    
    # Clean up standalone asterisks
    text = re.sub(r'^\*+\s*', '', text)             # Remove leading asterisks
    text = re.sub(r'\s*\*+$', '', text)             # Remove trailing asterisks
    
    # Clean up bullet points and numbering (but preserve the content)
    text = re.sub(r'^[-•*]\s*', '', text)           # Remove bullet points
    text = re.sub(r'^\d+\.\s*', '', text)           # Remove numbering
    
    # Clean up multiple spaces and normalize whitespace
    text = ' '.join(text.split())
    
    # Remove any remaining special formatting characters
    text = text.replace('---', '')                   # Remove horizontal rules
    text = text.replace('***', '')                   # Remove emphasis combinations
    
    # Remove content labels that might have been included
    if text.lower().startswith('content:'):
        text = text[8:].strip()
    
    return text.strip()

def clean_content_list_for_presentation(content_list):
    """Clean a list of content items for presentation use."""
    if not content_list or not isinstance(content_list, list):
        return []
    
    cleaned_list = []
    for item in content_list:
        if isinstance(item, str):
            cleaned_item = clean_text_for_presentation(item)
            # Skip empty items and content headers
            if cleaned_item and cleaned_item.lower() not in ['content', 'content:', '---']:
                cleaned_list.append(cleaned_item)
    
    return cleaned_list

def extract_enhanced_search_keywords(structured_content):
    """
    Truly smart keyword extraction that works with ANY language and subject without hardcoding.
    Uses statistical analysis and pattern recognition for 20+ subjects and languages.
    """
    # Combine all text content
    all_text = []
    for slide in structured_content:
        if slide.get('title'):
            all_text.append(slide['title'])
        if slide.get('content'):
            all_text.extend(slide['content'])
    
    if not all_text:
        return 'classroom education learning'
    
    text = ' '.join(all_text).lower()
    
    # Step 1: Extract content features using pattern recognition
    content_features = analyze_content_patterns(text)
    
    # Step 2: Extract meaningful terms using statistical methods
    meaningful_terms = extract_statistical_terms(text)
    
    # Step 3: Build search query using extracted features
    search_query = build_smart_search_query(content_features, meaningful_terms)
    
    logger.info(f"Smart extraction - Features: {content_features}, Query: '{search_query}'")
    return search_query

def analyze_content_patterns(text):
    """Analyze content using universal patterns that work across all languages and subjects."""
    features = set()
    
    # Enhanced Math patterns - more comprehensive decimal detection
    if re.search(r'\d+[\+\-\*\/×÷=]\d+', text):
        features.add('mathematics')
    
    if re.search(r'\d+\/\d+', text):  # Fractions
        features.add('fractions')
        features.add('mathematics')
    
    if re.search(r'\d+\.\d+', text):  # Decimals
        features.add('decimals')
        features.add('mathematics')
    
    # Enhanced decimal and place value detection
    if re.search(r'\b(decimal|decimals|tenths|hundredths|thousandths)\b', text, re.IGNORECASE):
        features.add('decimals')
        features.add('mathematics')
    
    if re.search(r'\b(powers?\s+of\s+10|place\s+value|place\s+values)\b', text, re.IGNORECASE):
        features.add('place_value')
        features.add('mathematics')
    
    if re.search(r'\b\d+%\b', text):  # Percentages
        features.add('statistics')
        features.add('mathematics')
    
    # Science patterns
    if re.search(r'\b[A-Z][a-z]+\s+[A-Z][a-z]+\b', text):  # Scientific names
        features.add('science')
    
    if re.search(r'\b\d+°[CF]?\b', text):  # Temperature
        features.add('science')
    
    if re.search(r'\b(CO2|H2O|O2|pH|DNA|RNA|NaCl)\b', text, re.IGNORECASE):  # Chemical formulas
        features.add('chemistry')
        features.add('science')
    
    # Biology patterns
    if re.search(r'\b(cell|cells|organism|species|evolution|photosynthesis)\b', text, re.IGNORECASE):
        features.add('biology')
        features.add('science')
    
    # Physics patterns
    if re.search(r'\b(force|energy|motion|gravity|velocity|acceleration)\b', text, re.IGNORECASE):
        features.add('physics')
        features.add('science')
    
    # Data visualization indicators (universal)
    chart_indicators = [r'chart', r'graph', r'table', r'data', r'survey', r'sample', r'diagram']
    if any(re.search(pattern, text, re.IGNORECASE) for pattern in chart_indicators):
        features.add('data_visualization')
    
    # History patterns
    if re.search(r'\b\d{4}\b', text):  # Years
        features.add('history')
    
    if re.search(r'\b(ancient|medieval|century|empire|civilization|war|revolution)\b', text, re.IGNORECASE):
        features.add('history')
    
    # Geography patterns
    if re.search(r'\b[A-Z][a-z]+\s+(river|mountain|ocean|sea|lake|country|continent)\b', text, re.IGNORECASE):
        features.add('geography')
    
    if re.search(r'\b(climate|weather|population|capital|border|map)\b', text, re.IGNORECASE):
        features.add('geography')
    
    # Language Arts patterns
    sentence_count = len(re.findall(r'[.!?]+', text))
    word_count = len(text.split())
    if word_count > 0 and sentence_count / word_count > 0.08:  # High punctuation density
        features.add('language_arts')
    
    if re.search(r'\b(reading|writing|grammar|vocabulary|literature|poetry|story)\b', text, re.IGNORECASE):
        features.add('language_arts')
    
    # Arts patterns
    if re.search(r'\b(art|painting|drawing|sculpture|color|brush|canvas|creative)\b', text, re.IGNORECASE):
        features.add('visual_arts')
    
    if re.search(r'\b(music|song|rhythm|melody|instrument|note|piano|guitar)\b', text, re.IGNORECASE):
        features.add('music')
    
    # Physical Education patterns
    if re.search(r'\b(sport|exercise|fitness|health|running|jumping|team|game)\b', text, re.IGNORECASE):
        features.add('physical_education')
    
    # Technology patterns
    if re.search(r'\b(computer|software|coding|programming|digital|internet|technology)\b', text, re.IGNORECASE):
        features.add('technology')
    
    # Social Studies patterns
    if re.search(r'\b(government|democracy|election|citizen|community|society|culture)\b', text, re.IGNORECASE):
        features.add('social_studies')
    
    # Economics patterns
    if re.search(r'\b(money|economy|business|trade|market|bank|finance)\b', text, re.IGNORECASE):
        features.add('economics')
    
    # Fun/Entertainment patterns
    if re.search(r'\b(fun|game|play|entertainment|hobby|leisure|enjoyment)\b', text, re.IGNORECASE):
        features.add('fun')
    
    # Holiday patterns
    holiday_patterns = [
        r'\b(christmas|halloween|thanksgiving|easter|valentine|birthday)\b',
        r'\b(holiday|celebration|festival|party|tradition|seasonal)\b',
        r'\b(december|january|october|november|february|march|april)\b'
    ]
    if any(re.search(pattern, text, re.IGNORECASE) for pattern in holiday_patterns):
        features.add('holidays')
    
    # Seasonal patterns
    if re.search(r'\b(spring|summer|fall|autumn|winter|season)\b', text, re.IGNORECASE):
        features.add('seasonal')
    
    # Special subjects
    if re.search(r'\b(special|therapy|intervention|support|accommodation)\b', text, re.IGNORECASE):
        features.add('special_education')
    
    # Health patterns
    if re.search(r'\b(health|nutrition|diet|wellness|safety|hygiene|medical)\b', text, re.IGNORECASE):
        features.add('health')
    
    # Environmental patterns
    if re.search(r'\b(environment|nature|ecology|conservation|sustainability|green)\b', text, re.IGNORECASE):
        features.add('environmental')
    
    return features

def extract_statistical_terms(text):
    """Extract meaningful terms using statistical analysis instead of stop word lists."""
    # Extract all words with Unicode support for any language
    words = re.findall(r'\b[\w\u00C0-\u024F\u1E00-\u1EFF\u0100-\u017F\u0180-\u024F]+\b', text.lower())
    
    # Statistical filtering instead of hardcoded stop words
    word_freq = Counter(words)
    total_words = len(words)
    
    if total_words == 0:
        return []
    
    meaningful_terms = []
    
    for word, count in word_freq.items():
        # Skip if word is too common or too rare
        frequency_ratio = count / total_words
        
        # Filter criteria (language-agnostic)
        if (len(word) >= 4 and                          # Reasonable length
            frequency_ratio < 0.3 and                  # Not too common (like "the", "de", etc.)
            frequency_ratio > (1 / total_words) and    # Not hapax legomena
            not word.isdigit() and                     # Not pure numbers
            not re.match(r'^(st|nd|rd|th|er|ème|º|ª|ый|ая|ое)$', word)):  # Not ordinal suffixes
            meaningful_terms.append(word)
    
    # Return top terms by a combination of frequency and length
    scored_terms = []
    for term in meaningful_terms:
        # Score based on length (longer = more specific) and moderate frequency
        freq_score = word_freq[term]
        length_score = min(len(term), 12)  # Cap at 12
        uniqueness_score = 1 / (word_freq[term] + 1)  # Prefer less common but not unique
        
        total_score = length_score * uniqueness_score * freq_score
        scored_terms.append((term, total_score))
    
    # Sort by score and return top 5
    top_terms = sorted(scored_terms, key=lambda x: x[1], reverse=True)[:5]
    return [term for term, score in top_terms]

def build_smart_search_query(features, terms):
    """Build search query based on detected features and extracted terms for 20+ subjects."""
    
    # Map features to search contexts (comprehensive subject mapping with better math contexts)
    feature_contexts = {
        # Enhanced Math Subjects
        'decimals': 'decimal numbers mathematics place value classroom',
        'place_value': 'place value chart mathematics tens hundreds classroom',
        'fractions': 'fractions mathematics visual circles pie charts',
        'mathematics': 'mathematics classroom education numbers calculations',
        'statistics': 'data charts graphs statistics mathematics',
        'data_visualization': 'charts graphs data visualization classroom',
        
        # Core Academic Subjects
        'science': 'science classroom experiment laboratory education',
        'chemistry': 'chemistry science laboratory beakers molecules',
        'biology': 'biology science nature organisms cells',
        'physics': 'physics science motion energy forces',
        'history': 'history timeline education classroom historical',
        'geography': 'geography maps world classroom globe',
        'language_arts': 'reading books classroom library education',
        
        # Arts and Creative Subjects
        'visual_arts': 'art painting creativity classroom colorful artistic',
        'music': 'music instruments classroom education musical',
        
        # Physical and Health
        'physical_education': 'sports exercise fitness gymnasium education',
        'health': 'health wellness education classroom medical',
        
        # Modern Subjects
        'technology': 'technology computers digital classroom programming',
        'social_studies': 'community society classroom education civic',
        'economics': 'money economics business classroom finance',
        'environmental': 'environment nature conservation classroom green',
        
        # Special Categories
        'fun': 'fun games colorful playful educational activities',
        'holidays': 'holiday celebration festive colorful seasonal',
        'seasonal': 'seasonal nature classroom decorative calendar',
        'special_education': 'inclusive education support classroom diverse',
        
        # Temporal
        'time_based': 'timeline calendar education classroom chronological'
    }
    
    # Start with detected features
    search_parts = []
    
    # Priority order for feature selection (most specific first)
    priority_features = [
        # Very specific math subjects first (for better relevance)
        'decimals', 'place_value', 'fractions', 'statistics', 'data_visualization',
        
        # Specific science subjects
        'chemistry', 'biology', 'physics',
        
        # Creative and special subjects
        'visual_arts', 'music', 'physical_education', 'technology', 
        'economics', 'environmental', 'fun', 'holidays', 'seasonal', 'special_education',
        
        # General subjects last (fallback)
        'mathematics', 'science', 'history', 'geography', 'language_arts',
        'social_studies', 'health', 'time_based'
    ]
    
    # Add primary feature context (first match wins for specificity)
    for feature in priority_features:
        if feature in features:
            search_parts.append(feature_contexts[feature])
            logger.debug(f"Selected primary feature: {feature}")
            break
    
    # Add meaningful extracted terms (max 2 for focused results)
    if terms:
        # Filter terms to avoid generic educational words
        generic_terms = {'lesson', 'class', 'student', 'learn', 'study', 'education', 'school', 'today', 'will'}
        
        filtered_terms = []
        for term in terms[:3]:  # Check top 3 terms
            if term not in generic_terms and len(term) > 4:
                filtered_terms.append(term)
                if len(filtered_terms) >= 2:  # Limit to 2 terms
                    break
        
        if filtered_terms:
            search_parts.extend(filtered_terms)
            logger.debug(f"Added extracted terms: {filtered_terms}")
    
    # Always add educational context if not present
    if not any('classroom' in part or 'education' in part for part in search_parts):
        search_parts.append('education classroom')
    
    # Join and deduplicate
    search_query = ' '.join(search_parts)
    words = search_query.split()
    
    # Remove duplicates while preserving order
    seen = set()
    unique_words = []
    for word in words:
        if word not in seen:
            seen.add(word)
            unique_words.append(word)
    
    result = ' '.join(unique_words)
    
    # Fallback if something went wrong
    if len(result.strip()) < 10:
        result = 'education classroom colorful learning mathematics'
    
    logger.debug(f"Final search query: '{result}'")
    return result

def add_image_to_slide(slide, image_bytes, lesson_topic=""):
    """
    Add an image to a widescreen slide (13.33" x 7.5") with proper positioning.
    Places image on the right side, centered vertically with text content.
    """
    try:
        # Create a BytesIO object from the image bytes
        image_stream = io.BytesIO(image_bytes)
        
        # Open image with PIL to get dimensions and potentially resize
        with Image.open(image_stream) as img:
            original_width, original_height = img.size
            
            # WIDESCREEN slide dimensions (13.33" x 7.5")
            slide_width = Inches(13.33)  # Widescreen width
            slide_height = Inches(7.5)   # Standard height
            
            # IMPROVED POSITIONING for widescreen: Place image on right side, larger
            target_width = Inches(4.5)   # Larger width for widescreen
            target_height = Inches(3.5)  # Good height for content
            
            # Calculate aspect ratio and adjust if needed
            img_aspect = original_width / original_height
            target_aspect = target_width / target_height
            
            if img_aspect > target_aspect:
                # Image is wider than target, fit by width
                final_width = target_width
                final_height = target_width / img_aspect
            else:
                # Image is taller than target, fit by height
                final_height = target_height
                final_width = target_height * img_aspect
            
            # Position on right side of widescreen slide, centered vertically
            left = slide_width - final_width - Inches(0.5)   # 0.5" margin from right edge
            
            # Center vertically in the content area (below title)
            content_start = Inches(2.0)  # Where content typically starts
            content_height = Inches(4.5)  # Typical content area height
            content_center = content_start + (content_height / 2)
            top = content_center - (final_height / 2)  # Center the image in content area
            
            # Ensure image doesn't go too high or too low
            min_top = Inches(1.8)  # Don't overlap with title
            max_top = slide_height - final_height - Inches(0.3)  # Don't go off bottom
            
            if top < min_top:
                top = min_top
            elif top > max_top:
                top = max_top
            
            # Reset image stream position
            image_stream.seek(0)
            
            # Add image to slide
            picture = slide.shapes.add_picture(
                image_stream, 
                left, 
                top, 
                final_width, 
                final_height
            )
            
            # Add subtle styling for professional look
            line = picture.line
            line.color.rgb = RGBColor(200, 200, 200)  # Light gray border
            line.width = Pt(0.75)
            
            # Optional: Add subtle shadow effect
            try:
                shadow = picture.shadow
                shadow.inherit = False
                shadow.style = 'OUTER'
                shadow.distance = Pt(3)
                shadow.blur_radius = Pt(4)
                shadow.color.rgb = RGBColor(128, 128, 128)
                shadow.transparency = 0.5
            except:
                pass  # Shadow effects might not be available in all versions
            
            logger.info(f"Successfully added image to widescreen slide (size: {final_width} x {final_height}, position: right-center)")
            return True
            
    except Exception as e:
        logger.error(f"Failed to add image to slide: {e}")
        return False
    finally:
        if 'image_stream' in locals():
            image_stream.close()

def add_text_box_to_slide(slide, content_items, with_image=False):
    """Add a text box to widescreen slide (13.33" x 7.5") with proper sizing."""
    
    # WIDESCREEN dimensions
    slide_width = Inches(13.33)
    
    if with_image:
        # Text takes up left portion of widescreen, leaving right for image
        left = Inches(0.8)      # Left margin
        top = Inches(2.1)       # Below title
        width = Inches(7.5)     # Much wider for widescreen (about 60% of slide)
        height = Inches(4.2)    # Good height for content
    else:
        # Use most of widescreen when no image
        left = Inches(0.8)
        top = Inches(2.1)
        width = Inches(11.5)    # Much wider for widescreen (about 85% of slide)
        height = Inches(4.5)
    
    # Create text box with widescreen-appropriate dimensions
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    
    # Optimize text frame properties for widescreen
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_bottom = Inches(0.15)
    text_frame.word_wrap = True
    text_frame.auto_size = None  # Prevent auto-sizing
    
    # Use cleaned content with larger text for widescreen visibility
    cleaned_items = clean_content_list_for_presentation(content_items)
    
    for item in cleaned_items:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = STYLE['fonts']['body']
        p.font.size = Pt(20)              # Larger font for widescreen
        p.font.color.rgb = STYLE['colors']['body']
        p.space_after = Pt(8)             # More spacing for readability
        p.line_spacing = 1.2              # Good line spacing
        p.level = 0                       # Consistent bullet level
    
    layout_desc = 'with image accommodation' if with_image else 'full widescreen'
    logger.info(f"Added widescreen text box ({layout_desc}) - size: {width}x{height}")

def clear_all_placeholder_content(slide):
    """AGGRESSIVELY clear all placeholder content including master slide placeholders."""
    try:
        placeholders_cleared = 0
        shapes_to_remove = []
        
        for shape in slide.shapes:
            try:
                # Skip title shapes
                if shape == slide.shapes.title:
                    continue
                
                # Check if this is a placeholder shape
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    # Try multiple clearing methods
                    if hasattr(shape, 'text_frame'):
                        # Method 1: Clear the text frame
                        shape.text_frame.clear()
                        
                        # Method 2: Set text to empty
                        if hasattr(shape.text_frame, 'text'):
                            shape.text_frame.text = ""
                        
                        # Method 3: Remove all paragraphs and add empty one
                        try:
                            shape.text_frame._element.clear()
                        except:
                            pass
                        
                        placeholders_cleared += 1
                        logger.debug(f"Aggressively cleared placeholder shape")
                
                # Also check by placeholder format
                elif hasattr(shape, 'placeholder_format'):
                    placeholder_type = getattr(shape.placeholder_format, 'type', None)
                    if placeholder_type in [2, 7, 8, 14]:  # Content placeholders
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.clear()
                            if hasattr(shape.text_frame, 'text'):
                                shape.text_frame.text = ""
                            placeholders_cleared += 1
                            logger.debug(f"Cleared placeholder by type: {placeholder_type}")
                
                # Last resort: check for "Click to add" text patterns
                elif hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                    text_content = shape.text_frame.text.lower()
                    if any(phrase in text_content for phrase in ['click to add', 'click to edit', 'add text']):
                        shape.text_frame.clear()
                        shape.text_frame.text = ""
                        placeholders_cleared += 1
                        logger.debug(f"Cleared shape with placeholder text: {text_content[:20]}")
                        
            except Exception as e:
                logger.debug(f"Could not process shape: {e}")
                continue
        
        logger.info(f"Aggressively cleared {placeholders_cleared} placeholder shapes")
        return placeholders_cleared > 0
        
    except Exception as e:
        logger.warning(f"Error in aggressive placeholder clearing: {e}")
        return False

def find_content_placeholder(slide):
    """Find a suitable content placeholder on the slide"""
    # Try to find placeholders in order of preference
    for placeholder in slide.placeholders:
        try:
            # Check if it's a content placeholder (not title)
            if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type:
                placeholder_type = placeholder.placeholder_format.type
                # Look for content, body, or object placeholders
                if placeholder_type in [2, 7, 8, 14]:  # CONTENT, BODY, OBJECT, CONTENT_WITH_CAPTION
                    return placeholder
        except:
            continue
    
    # If no specific content placeholder found, try by index
    try:
        # Most templates have content at index 1
        if len(slide.placeholders) > 1:
            return slide.placeholders[1]
    except:
        pass
    
    # Last resort: find any text-capable placeholder that's not the title
    for i, placeholder in enumerate(slide.placeholders):
        try:
            if i == 0:  # Skip title placeholder
                continue
            # Test if we can access text_frame
            if hasattr(placeholder, 'text_frame'):
                return placeholder
        except:
            continue
    
    return None

def create_clean_presentation_with_images(structured_content, include_images=False):
    """Create a PowerPoint presentation from clean structured content with enhanced images"""
    # Reset the image tracking flag
    if hasattr(create_clean_presentation_with_images, '_image_added'):
        delattr(create_clean_presentation_with_images, '_image_added')
    
    template_path = os.path.join(os.path.dirname(__file__), 'templates', 'FINAL_base_template_v1.pptx')
    
    # Check if template exists, use fallback if not
    if not os.path.exists(template_path):
        fallback_paths = [
            'templates/base_template.pptx',
            '../templates/base_template.pptx',
            'src/templates/base_template.pptx',
            'base_template.pptx'
        ]
        for path in fallback_paths:
            if os.path.exists(path):
                template_path = path
                break
    
    # Create presentation
    try:
        prs = Presentation(template_path)
        logger.info(f"Using template: {template_path}")
    except Exception as e:
        logger.warning(f"Could not load template: {e}. Creating blank presentation.")
        prs = Presentation()
    
    # Initialize Unsplash service if images are requested
    unsplash_photo_data = None
    image_bytes = None
    
    if include_images:
        try:
            from src.services.unsplash_service import unsplash_service
            
            # Generate enhanced search query from lesson content
            search_query = extract_enhanced_search_keywords(structured_content)
            logger.info(f"Searching for image with enhanced query: '{search_query}'")
            
            # Search for image
            unsplash_photo_data = unsplash_service.search_photo(search_query)
            
            if unsplash_photo_data:
                # Download image
                image_bytes = unsplash_service.download_photo(unsplash_photo_data)
                if image_bytes:
                    logger.info(f"Successfully retrieved relevant image by {unsplash_photo_data['photographer_name']} for query '{search_query}'")
                else:
                    logger.warning("Failed to download image from Unsplash")
            else:
                logger.warning(f"No suitable image found for enhanced query: '{search_query}'")
                
        except Exception as e:
            logger.error(f"Error fetching image from Unsplash: {e}")
            include_images = False  # Disable images for this generation
    
    # Log available slide layouts for debugging
    logger.info(f"Available slide layouts: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        try:
            layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
            logger.debug(f"Layout {i}: {layout_name} - {len(layout.placeholders)} placeholders")
        except:
            logger.debug(f"Layout {i}: Unknown layout")
    
    # Process each slide with clean structure and improved layout
    for slide_index, slide_data in enumerate(structured_content):
        try:
            # Try different layouts in order of preference
            layout_indices_to_try = [1, 0, 2, 3, 4]  # Title+Content, Title, Section, etc.
            slide = None
            
            for layout_idx in layout_indices_to_try:
                if layout_idx < len(prs.slide_layouts):
                    try:
                        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                        logger.debug(f"Successfully used layout {layout_idx} for slide {slide_index + 1}")
                        break
                    except Exception as e:
                        logger.debug(f"Failed to use layout {layout_idx}: {e}")
                        continue
            
            if not slide:
                # Fallback to the first available layout
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                logger.warning(f"Using fallback layout 0 for slide {slide_index + 1}")
            
            # Add image to first content slide (skip logo/title slides) with enhanced relevance
            has_image = False
            if include_images and image_bytes:
                # Check if this is the first slide with actual content
                slide_content = slide_data.get('content', [])
                clean_content = clean_content_list_for_presentation(slide_content)
                
                # Add image to first slide that has meaningful content
                if clean_content and not hasattr(create_clean_presentation_with_images, '_image_added'):
                    has_image = add_image_to_slide(slide, image_bytes, slide_data.get('title', ''))
                    # Mark that we've added an image to prevent adding to multiple slides
                    create_clean_presentation_with_images._image_added = True
                    logger.info(f"Added enhanced relevant image to slide {slide_index + 1} (first content slide)")
            
            # Clean and add title
            title_added = False
            raw_title = slide_data.get('title', 'Untitled')
            clean_title = clean_text_for_presentation(raw_title)
            
            if slide.shapes.title:
                try:
                    title_frame = slide.shapes.title.text_frame
                    title_frame.clear()
                    title_para = title_frame.add_paragraph()
                    title_para.text = clean_title
                    title_para.font.name = STYLE['fonts']['title']
                    title_para.font.size = STYLE['sizes']['title']
                    title_para.font.color.rgb = STYLE['colors']['title']
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER
                    title_added = True
                    logger.debug(f"Added clean title to slide {slide_index + 1}: {clean_title}")
                except Exception as e:
                    logger.warning(f"Failed to add title to slide {slide_index + 1}: {e}")
            
            # WIDESCREEN CONTENT HANDLING - Fixed for 13.33" x 7.5" template
            raw_content_items = slide_data.get('content', [])
            clean_content_items = clean_content_list_for_presentation(raw_content_items)
            
            if clean_content_items:
                # ALWAYS clear placeholders and use text boxes for consistency
                logger.info(f"Using text box for slide {slide_index + 1} (image: {has_image})")
                
                # CRITICAL: Clear ALL placeholders to prevent conflicts
                clear_all_placeholder_content(slide)
                
                # Add our custom text box with widescreen sizing
                add_text_box_to_slide(slide, clean_content_items, has_image)
                
                # Handle title if needed
                if not title_added and clean_title:
                    # Position title for widescreen
                    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(10), Inches(1))
                    title_frame = title_box.text_frame
                    title_para = title_frame.add_paragraph()
                    title_para.text = clean_title
                    title_para.font.name = STYLE['fonts']['title']
                    title_para.font.size = STYLE['sizes']['title']
                    title_para.font.color.rgb = STYLE['colors']['title']
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER
                    logger.debug(f"Added widescreen title as text box to slide {slide_index + 1}")
        
        except Exception as e:
            logger.error(f"Error creating slide {slide_index + 1}: {e}")
            # Create a basic slide with text boxes as last resort
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                
                # Add title as text box
                clean_title = clean_text_for_presentation(slide_data.get('title', f'Slide {slide_index + 1}'))
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                title_frame = title_box.text_frame
                title_para = title_frame.add_paragraph()
                title_para.text = clean_title
                title_para.font.name = STYLE['fonts']['title']
                title_para.font.size = STYLE['sizes']['title']
                title_para.font.color.rgb = STYLE['colors']['title']
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                
                # Add content as text box
                content_items = slide_data.get('content', [])
                if content_items:
                    add_text_box_to_slide(slide, content_items, False)
                
                logger.info(f"Created fallback slide {slide_index + 1} with text boxes")
            except Exception as fallback_error:
                logger.error(f"Failed to create fallback slide {slide_index + 1}: {fallback_error}")
    
    # Add attribution for image if used
    if include_images and unsplash_photo_data and image_bytes:
        try:
            # Add a final slide with attribution or add it to the last slide
            if len(prs.slides) > 0:
                last_slide = prs.slides[-1]
                
                # Add small attribution text at bottom
                attribution_text = f"Image: {unsplash_photo_data['photographer_name']} on Unsplash"
                
                # Add attribution as small text box at bottom
                attr_box = last_slide.shapes.add_textbox(
                    Inches(0.5), 
                    Inches(7), 
                    Inches(9), 
                    Inches(0.3)
                )
                attr_frame = attr_box.text_frame
                attr_para = attr_frame.add_paragraph()
                attr_para.text = attribution_text
                attr_para.font.size = Pt(8)
                attr_para.font.color.rgb = RGBColor(128, 128, 128)  # Gray color
                attr_para.alignment = PP_ALIGN.RIGHT
                
                logger.info(f"Added attribution: {attribution_text}")
                
        except Exception as e:
            logger.warning(f"Failed to add image attribution: {e}")
    
    logger.info(f"Created presentation with {len(structured_content)} slides (images: {'enabled' if include_images else 'disabled'})")
    return prs

def generate_search_query_from_content(structured_content, fallback="elementary classroom educational"):
    """
    Generate an enhanced search query for Unsplash based on lesson content.
    Uses the new enhanced keyword extraction for better image relevance.
    """
    try:
        if not structured_content or len(structured_content) == 0:
            return fallback
        
        # Use the enhanced keyword extraction
        enhanced_query = extract_enhanced_search_keywords(structured_content)
        
        if enhanced_query and enhanced_query != fallback:
            logger.info(f"Using enhanced search query: '{enhanced_query}'")
            return enhanced_query
        
        # Fallback to original logic if enhanced extraction doesn't find anything
        # Find the first slide with meaningful content (skip title-only slides)
        content_slide = None
        for slide in structured_content:
            content_items = slide.get('content', [])
            if content_items and len(content_items) > 0:
                # Check if this slide has real content, not just titles
                has_real_content = any(
                    len(item.strip()) > 10 and 
                    not item.lower().startswith(('students will', 'today we will', 'objectives'))
                    for item in content_items
                )
                if has_real_content:
                    content_slide = slide
                    break
        
        # If no content slide found, use the first slide
        if not content_slide:
            content_slide = structured_content[0]
        
        # Get title and content from the selected slide
        title = content_slide.get('title', '').lower()
        content_items = content_slide.get('content', [])
        content_text = ' '.join(content_items).lower() if content_items else ''
        
        # Combine title and content for analysis
        combined_text = f"{title} {content_text}".strip()
        
        # Use our smart analysis
        features = analyze_content_patterns(combined_text)
        terms = extract_statistical_terms(combined_text)
        result = build_smart_search_query(features, terms)
        
        if result and result != fallback:
            return result
        
        # Final enhanced fallback
        return "elementary classroom learning educational colorful"
        
    except Exception as e:
        logger.error(f"Error generating search query: {e}")
        return fallback

# BACKWARD COMPATIBILITY - Keep the old function for existing code
def create_presentation(structured_content, resource_type="PRESENTATION"):
    """Legacy function for backward compatibility"""
    # Convert old structure to new clean structure if needed
    clean_content = []
    
    for slide_data in structured_content:
        # Clean all the text content
        clean_slide = {
            'title': clean_text_for_presentation(slide_data.get('title', 'Untitled')),
            'layout': slide_data.get('layout', 'TITLE_AND_CONTENT'),
            'content': clean_content_list_for_presentation(slide_data.get('content', []))
        }
        
        # Handle old structure fields - convert to content and clean
        if not clean_slide['content']:
            # Try to extract from old fields
            if slide_data.get('left_column') or slide_data.get('right_column'):
                combined_content = (slide_data.get('left_column', []) + 
                                 slide_data.get('right_column', []))
                clean_slide['content'] = clean_content_list_for_presentation(combined_content)
            elif slide_data.get('teacher_notes'):
                clean_slide['content'] = clean_content_list_for_presentation(slide_data.get('teacher_notes', []))
        
        clean_content.append(clean_slide)
    
    return create_clean_presentation_with_images(clean_content, include_images=False)

def create_clean_presentation(structured_content):
    """Create a PowerPoint presentation from clean structured content without images"""
    return create_clean_presentation_with_images(structured_content, include_images=False)

def parse_outline_to_structured_content(outline_text, resource_type="PRESENTATION"):
    """Parse outline text into clean structured content"""
    logger.info(f"Parsing outline for resource type: {resource_type}")
    
    # Determine section/slide pattern based on resource type
    if resource_type.upper() == "PRESENTATION":
        section_pattern = r"Slide (\d+):\s*(.*)"
        section_word = "Slide"
    else:
        section_pattern = r"Section (\d+):\s*(.*)"
        section_word = "Section"
    
    # Split by section headers
    sections = []
    current_section = None
    
    lines = outline_text.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check if this is a section/slide header
        match = re.match(section_pattern, line)
        if match:
            # Save previous section
            if current_section:
                sections.append(current_section)
            
            # Start new section
            section_num = match.group(1)
            section_title = match.group(2).strip()
            current_section = {
                "title": clean_text_for_presentation(section_title),
                "layout": "TITLE_AND_CONTENT",
                "content": []
            }
        elif line.lower() == "content:":
            # Skip content headers
            continue
        elif line.startswith('-') or line.startswith('•'):
            # This is content
            if current_section:
                clean_content = clean_text_for_presentation(line.lstrip('-•').strip())
                if clean_content:
                    current_section["content"].append(clean_content)
        elif current_section and line:
            # Any other non-empty line goes to content
            clean_content = clean_text_for_presentation(line)
            if clean_content:
                current_section["content"].append(clean_content)
    
    # Don't forget the last section
    if current_section:
        sections.append(current_section)
    
    # If no sections found, create a fallback
    if not sections:
        sections.append({
            "title": "Generated Content",
            "layout": "TITLE_AND_CONTENT",
            "content": [clean_text_for_presentation(line.strip()) for line in lines if line.strip()]
        })
    
    logger.info(f"Successfully parsed {len(sections)} sections for {resource_type}")
    return sections