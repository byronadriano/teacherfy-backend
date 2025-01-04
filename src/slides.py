from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def parse_outline_to_structured_content(outline_text):
    
    """Parse the outline text into structured slide content with notes"""
    slides = []
    current_slide = None
    current_section = None
    
    lines = outline_text.strip().split('\n')
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        if line.startswith('Slide '):
            if current_slide:
                slides.append(current_slide)
            
            # Extract slide number and title
            title = line.split(':', 1)[1].strip() if ':' in line else line
            
            # Determine layout based on content
            layout = "TWO_COLUMN" if any(word in title.lower() for word in 
                ["vs", "comparison", "contrast"]) else "TITLE_AND_CONTENT"
            
            current_slide = {
                'title': title,
                'layout': layout,
                'content': [],
                'teacher_notes': [],
                'visual_elements': [],
                'left_column': [],
                'right_column': []
            }
            current_section = 'content'
            
        elif line.lower().startswith('content:'):
            current_section = 'content'
        elif line.lower().startswith('teacher notes:'):
            current_section = 'teacher_notes'
        elif line.lower().startswith('visual elements:'):
            current_section = 'visual_elements'
        elif line.startswith(('-', '*', '•')):
            content = line.lstrip('-*• ').strip()
            if current_slide:
                if current_slide['layout'] == "TWO_COLUMN" and current_section == 'content':
                    if len(current_slide['left_column']) <= len(current_slide['right_column']):
                        current_slide['left_column'].append(content)
                    else:
                        current_slide['right_column'].append(content)
                else:
                    current_slide[current_section].append(content)
        else:
            # Handle non-bullet point content
            if current_slide and not line.lower().endswith(':'):
                current_slide[current_section].append(line)
    
    if current_slide:
        slides.append(current_slide)
    
    return slides

def create_presentation(outline_json):
    """Create a PowerPoint presentation with notes from structured content"""
    # Load the template
    template_path = os.path.join(os.path.dirname(__file__), 'templates', 'base_template.pptx')
    prs = Presentation(template_path)
    
    # Title slide is already included in template (at index 0)
    
    for slide_data in outline_json:
        # Choose layout based on slide type
        layout_idx = 1  # Default to title and content (1)
        if slide_data['layout'] == "TWO_COLUMN":
            layout_idx = 3  # Two content layout (3)
        
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Add title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = slide_data['title']
            # Format title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                paragraph.font.name = 'Calibri'
        
        if slide_data['layout'] == "TWO_COLUMN":
            # Get placeholders for two-column layout
            shapes = [shape for shape in slide.placeholders]
            if len(shapes) >= 3:  # Title + 2 content placeholders
                left = shapes[1]
                right = shapes[2]
                
                # Add content to left column
                if slide_data['left_column']:
                    text_frame = left.text_frame
                    text_frame.clear()
                    for item in slide_data['left_column']:
                        p = text_frame.add_paragraph()
                        p.text = item
                        p.font.size = Pt(18)
                        p.font.name = 'Calibri'
                        p.level = 1
                
                # Add content to right column
                if slide_data['right_column']:
                    text_frame = right.text_frame
                    text_frame.clear()
                    for item in slide_data['right_column']:
                        p = text_frame.add_paragraph()
                        p.text = item
                        p.font.size = Pt(18)
                        p.font.name = 'Calibri'
                        p.level = 1
        else:
            # Add content to single-column layout
            shapes = [shape for shape in slide.placeholders]
            if len(shapes) >= 2:  # Title + content placeholder
                content_placeholder = shapes[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for item in slide_data['content']:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(18)
                    p.font.name = 'Calibri'
                    # Add bullet points for list items
                    if item.lstrip().startswith(('-', '•', '*')):
                        p.level = 1
                        p.text = item.lstrip('-•* ').strip()
        
        # Add notes
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame
        
        notes_text.clear()
        
        # Add teacher notes
        if slide_data['teacher_notes']:
            notes_text.text = "TEACHER NOTES:\n"
            for note in slide_data['teacher_notes']:
                p = notes_text.add_paragraph()
                p.text = f"• {note}"
        
        # Add visual elements
        if slide_data['visual_elements']:
            if notes_text.text:
                p = notes_text.add_paragraph()
                p.text = "\nVISUAL ELEMENTS:"
            else:
                notes_text.text = "VISUAL ELEMENTS:\n"
            
            for visual in slide_data['visual_elements']:
                p = notes_text.add_paragraph()
                p.text = f"• {visual}"
    
    return prs